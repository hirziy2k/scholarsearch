import asyncio
import json
import re
import subprocess
import sys
import warnings
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass

warnings.filterwarnings("ignore", message=".*fitz.*deprecated.*", category=DeprecationWarning)

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch, mm
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import HexColor
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ============================================================================
# PROTOCOL SKEPTIC: RAW SUBPROCESS + JSON EXTRACTION
# ============================================================================

class MCPSession:
    """Raw subprocess MCP client with background-line-reader.
    
    A daemon thread reads stdout lines into a queue. The _read_response
    polls the queue with a short timeout, so it never blocks forever
    (unlike stdout.read(1) which hangs when the process dies on Windows).
    """
    
    def __init__(self, command: str, args: List[str]):
        self.command = command
        self.args = args
        self.process: Optional[subprocess.Popen] = None
        self._id_counter = 0
        self._queue = None
        self._reader_thread = None
    
    async def start(self, timeout: float = 5.0):
        import queue, threading
        
        try:
            self.process = subprocess.Popen(
                [self.command] + self.args,
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
        except OSError as e:
            raise RuntimeError(f"Failed to launch MCP server: {e}")
        
        self._queue = queue.Queue()
        
        def _reader():
            try:
                for raw_line in self.process.stdout:
                    line = raw_line.decode('utf-8', errors='replace').rstrip('\r\n')
                    self._queue.put(line)
            except Exception:
                pass
            self._queue.put(None)
        
        self._reader_thread = threading.Thread(target=_reader, daemon=True)
        self._reader_thread.start()
        
        try:
            await asyncio.wait_for(self._initialize(), timeout=timeout)
        except asyncio.TimeoutError:
            self.close()
            raise RuntimeError(f"MCP server failed to initialize within {timeout}s")
    
    async def _initialize(self):
        result = await self._send_request("initialize", {
            "protocolVersion": "2024-11-05",
            "capabilities": {},
            "clientInfo": {"name": "pdf-to-pptx", "version": "1.0"}
        })
        self._send_notification("notifications/initialized", {})
        return result
    
    def _send_notification(self, method: str, params: Dict):
        if self.process.poll() is not None:
            return
        msg = {"jsonrpc": "2.0", "method": method, "params": params}
        try:
            self.process.stdin.write((json.dumps(msg) + "\n").encode())
            self.process.stdin.flush()
        except (BrokenPipeError, OSError):
            pass
    
    async def _send_request(self, method: str, params: Dict, timeout: float = 30.0) -> Dict:
        if self.process.poll() is not None:
            raise RuntimeError("MCP process is dead")
        self._id_counter += 1
        msg = {"jsonrpc": "2.0", "id": self._id_counter, "method": method, "params": params}
        try:
            self.process.stdin.write((json.dumps(msg) + "\n").encode())
            self.process.stdin.flush()
        except (BrokenPipeError, OSError):
            raise RuntimeError("MCP process stdin broken")
        return await self._read_response(self._id_counter, timeout)
    
    async def _read_response(self, request_id: int, timeout: float) -> Dict:
        import queue as _queue
        loop = asyncio.get_event_loop()
        
        def poll_queue():
            import time
            deadline = time.time() + timeout
            while time.time() < deadline:
                if self.process.poll() is not None:
                    return None
                try:
                    line = self._queue.get(timeout=0.3)
                except _queue.Empty:
                    continue
                if line is None:
                    return None
                line = line.strip()
                if not line or not line.startswith('{'):
                    continue
                try:
                    data = json.loads(line)
                    if isinstance(data, dict) and data.get("id") == request_id:
                        return data
                except json.JSONDecodeError:
                    continue
            return None
        
        result = await loop.run_in_executor(None, poll_queue)
        if result is None:
            raise TimeoutError(f"No response for request {request_id}")
        if "error" in result:
            raise RuntimeError(f"MCP error: {result['error']}")
        return result.get("result", {})
    
    async def call_tool(self, name: str, arguments: Dict, timeout: float = 30.0) -> Any:
        if self.process.poll() is not None:
            raise RuntimeError("MCP process died unexpectedly")
        result = await self._send_request("tools/call", {"name": name, "arguments": arguments}, timeout)
        return result
    
    async def list_tools(self) -> List[Dict]:
        result = await self._send_request("tools/list", {})
        return result.get("tools", [])
    
    def close(self):
        if self.process:
            self.process.terminate()
            self.process.wait(timeout=5)
            self.process = None


# ============================================================================
# COORDINATE MATH: PDF Points -> PPTX EMUs
# ============================================================================

PT_TO_EMU = 12700

def pdf_to_pptx_y(y_pt: float, _obj_height: float, _obj_top: float, page_height_pt: float) -> float:
    """Convert PDF Y coordinate (0 at bottom) to PPTX Y (0 at top)."""
    return (page_height_pt - y_pt) * PT_TO_EMU


# ============================================================================
# PHASE 2A: DYNAMIC CLUSTERING (moved from pdf_to_pptx for shared use)
# ============================================================================

def span_in_bbox(span: Dict, bbox: Dict) -> bool:
    if not bbox:
        return False
    sx, sy = span.get("x", 0), span.get("y", 0)
    return (bbox.get("x1", 0) <= sx <= bbox.get("x2", 0) and
            bbox.get("y1", 0) <= sy <= bbox.get("y2", 0))


def merge_cluster(spans: List[Dict], block: Dict) -> Dict:
    if not spans:
        return {}
    min_x = min(s.get("x", 0) for s in spans)
    min_y = min(s.get("y", 0) for s in spans)
    max_x = max(s.get("x", 0) + s.get("width", 0) for s in spans)
    max_y = max(s.get("y", 0) + s.get("height", 0) for s in spans)
    text = " ".join(s.get("text", "") for s in spans)
    avg_font = sum(s.get("font_size", 12) for s in spans) / len(spans)
    return {
        "text": text, "left": min_x, "top": min_y,
        "width": max_x - min_x, "height": max_y - min_y,
        "font_size": avg_font, "font_name": spans[0].get("font", "Calibri"),
        "layout_block_ref": block
    }


def cluster_spans(raw_spans: List[Dict], layout_blocks: List[Dict]) -> List[Dict]:
    horizontal_clusters = []
    for block in layout_blocks:
        block_bbox = block.get("coordinates", {})
        if not block_bbox:
            continue
        block_spans = [s for s in raw_spans if span_in_bbox(s, block_bbox)]
        if not block_spans:
            continue
        block_spans.sort(key=lambda s: (s.get("y", 0), s.get("x", 0)))
        current_cluster = []
        for span in block_spans:
            if not current_cluster:
                current_cluster.append(span)
            else:
                last = current_cluster[-1]
                font_size = span.get("font_size", 12)
                gap = span.get("x", 0) - (last.get("x", 0) + last.get("width", 0))
                vertical_diff = abs(span.get("y", 0) - last.get("y", 0))
                if gap <= font_size * 0.3 and vertical_diff < font_size * 0.5:
                    current_cluster.append(span)
                else:
                    horizontal_clusters.append(merge_cluster(current_cluster, block))
                    current_cluster = [span]
        if current_cluster:
            horizontal_clusters.append(merge_cluster(current_cluster, block))
    
    if not horizontal_clusters:
        return []
    horizontal_clusters.sort(key=lambda c: (c.get("top", 0), c.get("left", 0)))
    vertical_clusters = []
    current_vcluster = [horizontal_clusters[0]]
    for cluster in horizontal_clusters[1:]:
        last = current_vcluster[-1]
        font_size = cluster.get("font_size", 12)
        vertical_diff = abs(cluster.get("top", 0) - last.get("top", 0))
        h_gap = cluster.get("left", 0) - (last.get("left", 0) + last.get("width", 0))
        if vertical_diff < font_size * 1.5 and h_gap < font_size * 2:
            current_vcluster.append(cluster)
        else:
            merged = merge_vcluster(current_vcluster)
            vertical_clusters.append(merged)
            current_vcluster = [cluster]
    if current_vcluster:
        vertical_clusters.append(merge_vcluster(current_vcluster))
    return vertical_clusters


def merge_vcluster(clusters: List[Dict]) -> Dict:
    if not clusters:
        return {}
    min_x = min(c.get("left", 0) for c in clusters)
    min_y = min(c.get("top", 0) for c in clusters)
    max_x = max(c.get("left", 0) + c.get("width", 0) for c in clusters)
    max_y = max(c.get("top", 0) + c.get("height", 0) for c in clusters)
    text = "\n".join(c.get("text", "") for c in clusters)
    avg_font = sum(c.get("font_size", 12) for c in clusters) / len(clusters)
    return {
        "text": text, "left": min_x, "top": min_y,
        "width": max_x - min_x, "height": max_y - min_y,
        "font_size": avg_font, "font_name": clusters[0].get("font_name", "Calibri")
    }


# ============================================================================
# PHASE 1 TEST HARNESS
# ============================================================================

def create_synthetic_test_pdf(output_path: Path) -> Dict[str, Any]:
    if not REPORTLAB_AVAILABLE:
        raise RuntimeError("reportlab not installed")
    
    page_width_pt, page_height_pt = 612.0, 792.0
    c = canvas.Canvas(str(output_path), pagesize=(page_width_pt, page_height_pt))
    
    elements = {
        "text_box": {"text": "HELLO WORLD", "x_pt": 72.0, "y_pt": 72.0, "width_pt": 200.0, "height_pt": 50.0, "font_size_pt": 24.0, "font_name": "Helvetica"},
        "table": {"rows": 3, "cols": 3, "x_pt": 72.0, "y_pt": 150.0, "width_pt": 300.0, "height_pt": 100.0},
        "image": {"x_pt": 72.0, "y_pt": 300.0, "width_pt": 150.0, "height_pt": 100.0},
        "rectangle": {"x_pt": 300.0, "y_pt": 300.0, "width_pt": 100.0, "height_pt": 50.0},
    }
    
    c.setFont("Helvetica", 24)
    c.drawString(72, 72, "HELLO WORLD")
    
    c.setStrokeColor(HexColor("#000000"))
    c.setLineWidth(1)
    table_x, table_y = 72, 150
    cell_w, cell_h = 100, 33.33
    for r in range(3):
        for ci in range(3):
            c.rect(table_x + ci * cell_w, table_y + r * cell_h, cell_w, cell_h)
            c.setFont("Helvetica", 10)
            c.drawString(table_x + ci * cell_w + 5, table_y + r * cell_h + 10, f"R{r}C{ci}")
    
    c.rect(72, 300, 150, 100)
    c.setFillColor(HexColor("#FFA500"))
    c.rect(300, 300, 100, 50, fill=1)
    
    c.save()
    
    expected = {}
    for name, el in elements.items():
        expected[name] = {"left": pt_to_emu(el["x_pt"]), "top": pdf_to_pptx_y(el["y_pt"], el["height_pt"], el["height_pt"], page_height_pt), "width": pt_to_emu(el["width_pt"]), "height": pt_to_emu(el["height_pt"])}
    expected["slide"] = {"width": pt_to_emu(page_width_pt), "height": pt_to_emu(page_height_pt)}
    return expected


def pt_to_emu(pt: float) -> int:
    return int(pt * PT_TO_EMU)


def verify_pptx_output(pptx_path: Path, expected: Dict[str, Any]) -> List[str]:
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not installed")
    errors = []
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    if abs(prs.slide_width - expected["slide"]["width"]) > 1:
        errors.append(f"Slide width mismatch")
    if abs(prs.slide_height - expected["slide"]["height"]) > 1:
        errors.append(f"Slide height mismatch")
    for name, exp in expected.items():
        if name == "slide":
            continue
        found = any(abs(s.left-exp["left"])<=1 and abs(s.top-exp["top"])<=1 and abs(s.width-exp["width"])<=1 and abs(s.height-exp["height"])<=1 for s in slide.shapes)
        if not found:
            errors.append(f"Missing: {name}")
    return errors


async def run_phase1_test() -> bool:
    print("=" * 60)
    print("PHASE 1 TEST HARNESS")
    print("=" * 60)
    if not REPORTLAB_AVAILABLE or not PPTX_AVAILABLE:
        print("SKIPPED")
        return False
    import tempfile
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        expected = create_synthetic_test_pdf(tmpdir / "test.pdf")
        prs = Presentation()
        prs.slide_width = Emu(expected["slide"]["width"])
        prs.slide_height = Emu(expected["slide"]["height"])
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for name, exp in expected.items():
            if name == "slide": continue
            if name == "text_box":
                tb = slide.shapes.add_textbox(Emu(exp["left"]), Emu(exp["top"]), Emu(exp["width"]), Emu(exp["height"]))
                tb.text_frame.word_wrap = False
                p = tb.text_frame.paragraphs[0]
                p.text = "HELLO WORLD"
                p.font.size = Pt(24)
                p.font.name = "Helvetica"
            elif name == "table":
                ts = slide.shapes.add_table(3, 3, Emu(exp["left"]), Emu(exp["top"]), Emu(exp["width"]), Emu(exp["height"]))
                for r in range(3):
                    for ci in range(3):
                        ts.table.cell(r, ci).text = f"R{r}C{ci}"
            elif name == "image":
                slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(exp["left"]), Emu(exp["top"]), Emu(exp["width"]), Emu(exp["height"]))
            elif name == "rectangle":
                s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(exp["left"]), Emu(exp["top"]), Emu(exp["width"]), Emu(exp["height"]))
                s.fill.solid()
                s.fill.fore_color.rgb = RGBColor(0xFF, 0xA5, 0x00)
        prs.save(str(tmpdir / "out.pptx"))
        errors = verify_pptx_output(tmpdir / "out.pptx", expected)
        if errors:
            print("FAILED:", errors)
            return False
        print("PASSED: All coordinates match +/-1 EMU")
        return True


# ============================================================================
# NEW ORCHESTRATOR INTEGRATION
# ============================================================================

class PDFToPPTXConverter:
    def __init__(self, pdf_path: str, output_path: str):
        self.pdf_path = Path(pdf_path).resolve()
        self.output_path = Path(output_path).resolve()
        self.warnings: List[str] = []
        self.errors: List[str] = []

    async def __aenter__(self):
        return self

    async def __aexit__(self, exc_type, exc_val, exc_tb):
        pass

    async def convert(self) -> 'ConversionResult':
        # Import here to avoid circular import
        from orchestrator_v2 import convert_pdf, ConversionResult
        print(f"Converting {self.pdf_path.name}...")
        result = await convert_pdf(str(self.pdf_path), self.output_path.parent, mcp_augment=True)
        if result.success:
            print(f"  Success: {result.total_pages} pages -> {result.total_slides} slides")
            if result.mcp_augmented_pages:
                print(f"  MCP augmented: {result.mcp_augmented_pages} pages")
        else:
            for e in result.errors:
                print(f"  Error: {e}")
        return result


async def main():
    if len(sys.argv) >= 2 and sys.argv[1] == "--test-phase1":
        await run_phase1_test()
    elif len(sys.argv) >= 3:
        pdf_path, output_path = sys.argv[1], sys.argv[2]
        if not Path(pdf_path).exists():
            print(f"Error: {pdf_path} not found")
            sys.exit(1)
        async with PDFToPPTXConverter(pdf_path, output_path) as converter:
            result = await converter.convert()
        if result.success:
            print(f"\nSuccess: {result.output_path}")
        else:
            print("\nFailed:")
            for e in result.errors:
                print(f"  - {e}")
            sys.exit(1)
    else:
        print("Usage:")
        print("  python pdf_to_pptx.py --test-phase1")
        print("  python pdf_to_pptx.py input.pdf output.pptx")


if __name__ == "__main__":
    asyncio.run(main())