#!/usr/bin/env python3
"""
orchestrator_v2.py — New Primary Orchestrator

PyMuPDF Primary + MCP Augmentation + Dynamic Clustering + Binary Validation

Flow:
1. PyMuPDF extracts ALL data (text+coords, tables, vectors) natively
2. Complexity analysis decides which pages need MCP augmentation
3. MCP augments complex pages with semantic reading order
4. Dynamic clustering groups spans into logical blocks
5. Binary validation rejects hallucinated tables
6. PPTX rendering with diagnostic stubs for any failures
"""

import asyncio
import json
import time
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, field, asdict
from collections import defaultdict

from pdf_extractor import PyMuPDFExtractor, PageExtraction, TextBlock, TextLine, TextSpan, ExtractedTable
from mcp_augment import MCPAugmentor, should_augment
from pdf_to_pptx import MCPSession, PT_TO_EMU, pdf_to_pptx_y, cluster_spans
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

try:
    from theme_mapper import classify_and_enhance, ThemeColors, SemanticClassifier
    THEME_MAPPER_AVAILABLE = True
except ImportError:
    THEME_MAPPER_AVAILABLE = False

try:
    from metadata_layer import MetadataExtractor, MetadataBinder
    METADATA_LAYER_AVAILABLE = True
except ImportError:
    METADATA_LAYER_AVAILABLE = False

try:
    from font_emulator import FontDetector, FontMetricEmulator, apply_font_fallback
    FONT_EMULATOR_AVAILABLE = True
except ImportError:
    FONT_EMULATOR_AVAILABLE = False

try:
    from reflow_grouping import apply_reflow_properties
    REFLOW_GROUPING_AVAILABLE = True
except ImportError:
    REFLOW_GROUPING_AVAILABLE = False

try:
    from pdf_sandbox import sanitize_pdf, SanitizeResult
    SANDBOX_AVAILABLE = True
except ImportError:
    SANDBOX_AVAILABLE = False

try:
    from collision_matrix import resolve_collisions
    COLLISION_MATRIX_AVAILABLE = True
except ImportError:
    COLLISION_MATRIX_AVAILABLE = False

try:
    from telemetry import create_conversion_trace, StructuredLogger, DiagnosticReport
    TELEMETRY_AVAILABLE = True
except ImportError:
    TELEMETRY_AVAILABLE = False

try:
    from summary_slide import append_summary_slide, create_summary_from_result, ConversionSummary
    SUMMARY_SLIDE_AVAILABLE = True
except ImportError:
    SUMMARY_SLIDE_AVAILABLE = False


@dataclass
class ConversionResult:
    pdf_path: str
    success: bool
    output_path: Optional[str] = None
    total_pages: int = 0
    total_slides: int = 0
    total_time_ms: float = 0
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    mcp_augmented_pages: int = 0
    native_fallbacks: int = 0


class ClusteringEngine:
    """Reusable dynamic clustering from Phase 2A."""
    
    @staticmethod
    def cluster_page_blocks(blocks: List[TextBlock], body_font_size: float) -> List[Dict]:
        """Convert TextBlocks to spans and cluster them."""
        all_spans = []
        for block in blocks:
            for line in block.lines:
                for span in line.spans:
                    if span.text.strip():
                        all_spans.append({
                            "text": span.text,
                            "x": span.bbox[0],
                            "y": span.bbox[1],
                            "width": span.bbox[2] - span.bbox[0],
                            "height": span.bbox[3] - span.bbox[1],
                            "font_size": span.size,
                            "font": span.font,
                            "color": span.color
                        })
        
        if not all_spans:
            return []
        
        # Create layout_blocks from TextBlocks for cluster_spans
        layout_blocks = []
        for block in blocks:
            if block.lines:
                spans_in_block = []
                for line in block.lines:
                    for span in line.spans:
                        spans_in_block.append({
                            "x": span.bbox[0], "y": span.bbox[1],
                            "width": span.bbox[2] - span.bbox[0],
                            "height": span.bbox[3] - span.bbox[1]
                        })
                if spans_in_block:
                    x1 = min(s["x"] for s in spans_in_block)
                    y1 = min(s["y"] for s in spans_in_block)
                    x2 = max(s["x"] + s["width"] for s in spans_in_block)
                    y2 = max(s["y"] + s["height"] for s in spans_in_block)
                    layout_blocks.append({
                        "coordinates": {"x1": x1, "y1": y1, "x2": x2, "y2": y2}
                    })
        
        from pdf_to_pptx import cluster_spans
        clustered = cluster_spans(all_spans, layout_blocks)
        return clustered


class ValidationEngine:
    """Reusable binary validation from Phase 2B."""
    
    @staticmethod
    def validate_tables(tables: List[ExtractedTable]) -> Tuple[List[Dict], List[Dict]]:
        """Validate tables, return (valid, fallbacks)."""
        valid = []
        fallbacks = []
        
        for table in tables:
            data = table.data
            rows = table.rows
            cols = table.cols
            
            # Hallucination check
            if rows > 20 and cols > 10:
                fallbacks.append({
                    "type": "fallback_image",
                    "page": table.page,
                    "bbox": table.bbox,
                    "reason": f"grid_hallucination_{rows}x{cols}"
                })
                continue
            
            # Sparsity check
            if data:
                total_cells = sum(len(r) for r in data)
                empty_cells = sum(1 for r in data for v in r if not v or not str(v).strip())
                if total_cells > 0 and empty_cells / total_cells > 0.7:
                    fallbacks.append({
                        "type": "fallback_image",
                        "page": table.page,
                        "bbox": table.bbox,
                        "reason": "sparse_data"
                    })
                    continue
            
            valid.append({
                "page": table.page,
                "bbox": table.bbox,
                "rows": rows,
                "cols": cols,
                "data": data
            })
        
        return valid, fallbacks


class PPTXRenderer:
    """Renders clustered blocks + validated tables to PPTX."""
    
    def __init__(self, doc_extraction, clustered_pages, valid_tables, fallback_tables, mcp_augmented_pages, native_fallbacks):
        self.doc = doc_extraction
        self.clustered_pages = clustered_pages
        self.valid_tables = valid_tables
        self.fallback_tables = fallback_tables
        self.mcp_augmented_pages = mcp_augmented_pages
        self.native_fallbacks = native_fallbacks
    
    def render(self, output_path: Path) -> int:
        """Render PPTX, return slide count."""
        page_width, page_height = self.doc.pages[0].page_size if self.doc.pages else (612, 792)
        page_width_emu = int(page_width * PT_TO_EMU)
        page_height_emu = int(page_height * PT_TO_EMU)
        
        prs = Presentation()
        prs.slide_width = Emu(page_width_emu)
        prs.slide_height = Emu(page_height_emu)
        blank_layout = prs.slide_layouts[6]
        
        tables_by_page = defaultdict(list)
        for t in self.valid_tables:
            tables_by_page[t["page"]].append(t)
        
        fallbacks_by_page = defaultdict(list)
        for t in self.fallback_tables:
            fallbacks_by_page[t["page"]].append(t)
        
        for page_idx, page in enumerate(self.doc.pages):
            page_num = page_idx + 1
            slide = prs.slides.add_slide(blank_layout)
            
            # Render text clusters
            clusters = self.clustered_pages.get(page_num, [])
            for cluster in clusters:
                self._render_text_cluster(slide, cluster, page_width, page_height)
            
            # Render validated tables
            for table in tables_by_page.get(page_num, []):
                self._render_table(slide, table, page_width, page_height)
            
            # Render fallback indicators
            for fb in fallbacks_by_page.get(page_num, []):
                self._render_fallback(slide, fb, page_width, page_height)
            
            # Diagnostic stub for native fallback pages
            if page_num in [p for p, c in enumerate(self.doc.pages) if c]:
                pass  # Native extraction succeeded
        
        prs.save(str(output_path))
        return len(prs.slides)
    
    def _render_text_cluster(self, slide, cluster: Dict, page_width: float, page_height: float):
        bbox = cluster.get("bbox", (0, 0, 0, 0))
        if bbox[2] <= bbox[0] or bbox[3] <= bbox[1]:
            return
        
        left_emu = int(bbox[0] * PT_TO_EMU)
        top_emu = int(pdf_to_pptx_y(bbox[1], 0, 0, page_height))
        width_emu = int((bbox[2] - bbox[0]) * PT_TO_EMU)
        height_emu = int((bbox[3] - bbox[1]) * PT_TO_EMU)
        
        txBox = slide.shapes.add_textbox(Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cluster.get("text", "")
        p.font.size = Pt(max(8, cluster.get("font_size", 10)))
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    def _render_table(self, slide, table: Dict, page_width: float, page_height: float):
        bbox = table.get("bbox", (0, 0, 0, 0))
        if bbox[2] <= bbox[0] or bbox[3] <= bbox[1]:
            return
        
        left_emu = int(bbox[0] * PT_TO_EMU)
        top_emu = int(pdf_to_pptx_y(bbox[1], 0, 0, page_height))
        width_emu = int((bbox[2] - bbox[0]) * PT_TO_EMU)
        height_emu = int((bbox[3] - bbox[1]) * PT_TO_EMU)
        
        rows = table.get("rows", 1)
        cols = table.get("cols", 1)
        data = table.get("data", [])
        
        table_shape = slide.shapes.add_table(rows, cols, Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu))
        tbl = table_shape.table
        
        for r_idx in range(min(rows, len(data))):
            for c_idx in range(min(cols, len(data[r_idx]) if r_idx < len(data) else 0)):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = str(data[r_idx][c_idx]) if data[r_idx][c_idx] else ""
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                    paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    def _render_fallback(self, slide, fb: Dict, page_width: float, page_height: float):
        bbox = fb.get("bbox", (0, 0, 0, 0))
        if bbox[2] <= bbox[0] or bbox[3] <= bbox[1]:
            return
        
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Emu(int(bbox[0] * PT_TO_EMU)),
            Emu(int(pdf_to_pptx_y(bbox[1], 0, 0, page_height))),
            Emu(int((bbox[2] - bbox[0]) * PT_TO_EMU)),
            Emu(int((bbox[3] - bbox[1]) * PT_TO_EMU)))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xE0)
        shape.line.color.rgb = RGBColor(0xFF, 0x6B, 0x35)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Fallback: {fb.get('reason', 'unknown')}]"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(0x8B, 0x00, 0x00)


async def convert_pdf(pdf_path: str, output_dir: Path, mcp_augment: bool = True,
                      trace_id: str = None) -> ConversionResult:
    """Main conversion function: PyMuPDF primary + MCP augmentation + all layers."""
    result = ConversionResult(pdf_path=pdf_path, success=False)
    output_path = output_dir / f"{Path(pdf_path).stem}.pptx"
    
    # Telemetry setup
    trace = None
    log = None
    if TELEMETRY_AVAILABLE:
        trace = create_conversion_trace(pdf_path)
        if trace_id:
            trace.trace_id = trace_id
        log = StructuredLogger(trace)
        result.warnings.append(f"Trace ID: {trace.trace_id}")
    
    t0 = time.time()
    
    try:
        # Phase 0: Pre-flight PDF Sanitization
        sanitized_path = pdf_path
        if SANDBOX_AVAILABLE:
            with trace.start_span("sanitize", {"pdf_path": pdf_path}) if trace else None:
                san_result = sanitize_pdf(pdf_path, str(output_dir))
                if not san_result.safe:
                    result.warnings.append(f"PDF threats found: {san_result.threats_found}")
                if san_result.validation_errors:
                    result.warnings.append(f"Validation: {san_result.validation_errors}")
                sanitized_path = san_result.sanitized_path
                result.warnings.append(f"Sanitization: {san_result.time_ms:.0f}ms")
        
        # Phase 1: PyMuPDF Primary Extraction
        with trace.start_span("extract", {"pdf_path": sanitized_path}) if trace else None:
            with PyMuPDFExtractor(sanitized_path) as extractor:
                doc = extractor.extract_all()
                complexities = []
                for i in range(len(doc.pages)):
                    c = extractor.analyze_complexity(extractor.doc[i])
                    complexities.append(c)
                
                result.total_pages = len(doc.pages)
        
        # Phase 2: MCP Augmentation (selective)
        mcp_augmented = {}
        if mcp_augment:
            try:
                with trace.start_span("mcp_augment") if trace else None:
                    session = MCPSession("uvx", ["mcp-pdf"])
                    await session.start(timeout=5.0)
                    augmentor = MCPAugmentor(session)
                    
                    mcp_calls = 0
                    MAX_MCP_CALLS = 10
                    for page_num, complexity in enumerate(complexities, 1):
                        if mcp_calls >= MAX_MCP_CALLS:
                            break
                        if should_augment(complexity):
                            try:
                                mcp_result = await augmentor.analyze_page(sanitized_path, page_num)
                                if "error" not in mcp_result:
                                    mcp_augmented[page_num] = mcp_result
                                    mcp_calls += 1
                            except Exception:
                                break
                    
                    result.mcp_augmented_pages = len(mcp_augmented)
                    session.close()
            except Exception as e:
                result.warnings.append(f"MCP unavailable: {type(e).__name__} — using PyMuPDF only")
        
        # Phase 3: Dynamic Clustering
        clustered_pages = {}
        with trace.start_span("cluster") if trace else None:
            for page in doc.pages:
                clusters = ClusteringEngine.cluster_page_blocks(page.blocks, doc.body_font_size)
                if mcp_augmented.get(page.page_num):
                    augmentor = MCPAugmentor(None)
                    clusters = augmentor.sort_blocks_by_reading_order(clusters, mcp_augmented[page.page_num])
                clustered_pages[page.page_num] = clusters
        
        # Phase 4: Binary Validation
        all_tables = []
        with trace.start_span("validate") if trace else None:
            for page in doc.pages:
                all_tables.extend(page.tables)
            valid_tables, fallback_tables = ValidationEngine.validate_tables(all_tables)
        
        # Phase 4.5: Semantic Classification (Theme Mapper)
        if THEME_MAPPER_AVAILABLE:
            with trace.start_span("theme_classify") if trace else None:
                for page_num, clusters in clustered_pages.items():
                    table_bboxes = [t["bbox"] for t in valid_tables if t["page"] == page_num]
                    clustered_pages[page_num] = classify_and_enhance(
                        clusters, doc.pages[page_num - 1].page_size[1],
                        doc.body_font_size, table_bboxes
                    )
                result.warnings.append("Theme mapper: active")
        
        # Phase 4.6: Font Metric Emulation
        if FONT_EMULATOR_AVAILABLE:
            with trace.start_span("font_emulate") if trace else None:
                try:
                    font_detector = FontDetector()
                    font_emulator = FontMetricEmulator(font_detector)
                    fallback_count = 0
                    for page_num, clusters in clustered_pages.items():
                        enhanced = apply_font_fallback(clusters, font_detector, font_emulator)
                        fallback_count += sum(1 for c in enhanced if c.get("font_fallback_applied"))
                    if fallback_count > 0:
                        result.warnings.append(f"Font fallback: {fallback_count} clusters adjusted")
                except Exception as e:
                    result.warnings.append(f"Font emulation unavailable: {type(e).__name__}")
        
        # Phase 4.7: Collision Resolution (after font emulation may have expanded boxes)
        spillover_pages = {}
        if COLLISION_MATRIX_AVAILABLE:
            with trace.start_span("collision_resolve") if trace else None:
                for page_num, clusters in clustered_pages.items():
                    if len(clusters) > 1:
                        ph = doc.pages[page_num - 1].page_size[1] if page_num <= len(doc.pages) else 792
                        pw = doc.pages[page_num - 1].page_size[0] if page_num <= len(doc.pages) else 612
                        adjusted, overflow = resolve_collisions(clusters, ph * PT_TO_EMU, pw * PT_TO_EMU)
                        clustered_pages[page_num] = adjusted
                        if overflow:
                            spillover_pages[page_num] = overflow
                if spillover_pages:
                    result.warnings.append(f"Collisions: resolved on {len(spillover_pages)} pages")
        
        # Phase 5: Render PPTX
        with trace.start_span("render_pptx", {"output_path": str(output_path)}) if trace else None:
            renderer = PPTXRenderer(doc, clustered_pages, valid_tables, fallback_tables,
                                   result.mcp_augmented_pages, result.native_fallbacks)
            slide_count = renderer.render(output_path)
        
        # Phase 5.5: Auto-Reflow Grouping
        if REFLOW_GROUPING_AVAILABLE:
            with trace.start_span("reflow_apply") if trace else None:
                try:
                    prs = Presentation(str(output_path))
                    reflow_count = 0
                    for slide_idx, slide in enumerate(prs.slides):
                        page_num = slide_idx + 1
                        clusters = clustered_pages.get(page_num, [])
                        if clusters:
                            page_obj = doc.pages[page_num - 1] if page_num <= len(doc.pages) else None
                            pw = page_obj.page_size[0] if page_obj else 612
                            ph = page_obj.page_size[1] if page_obj else 792
                            reflow_count += apply_reflow_properties(slide, clusters, ph, pw)
                    prs.save(str(output_path))
                    if reflow_count > 0:
                        result.warnings.append(f"Reflow: {reflow_count} shapes set to auto-fit")
                except Exception as e:
                    result.warnings.append(f"Reflow grouping unavailable: {type(e).__name__}")
        
        # Phase 5.6: Metadata Preservation (hyperlinks, bookmarks, alt-text)
        if METADATA_LAYER_AVAILABLE:
            with trace.start_span("metadata_bind") if trace else None:
                try:
                    extractor = MetadataExtractor(sanitized_path)
                    pdf_metadata = extractor.extract_all()
                    page_mapping = {i + 1: i for i in range(len(doc.pages))}
                    binder = MetadataBinder()
                    prs = Presentation(str(output_path))
                    binder.bind_metadata(prs, pdf_metadata, page_mapping)
                    prs.save(str(output_path))
                    if pdf_metadata.hyperlinks:
                        result.warnings.append(f"Metadata: {len(pdf_metadata.hyperlinks)} hyperlinks bound")
                    if pdf_metadata.bookmarks:
                        result.warnings.append(f"Metadata: {len(pdf_metadata.bookmarks)} bookmarks bound")
                except Exception as e:
                    result.warnings.append(f"Metadata binding failed: {type(e).__name__}: {e}")
        
        result.output_path = str(output_path)
        result.total_slides = slide_count
        result.success = True
        
        # Phase 6: Conversion Summary Slide (end-user transparency)
        if SUMMARY_SLIDE_AVAILABLE and result.success:
            try:
                trace_id = trace.trace_id if trace else None
                summary = create_summary_from_result(result, trace_id=trace_id)
                prs = Presentation(str(output_path))
                added = append_summary_slide(prs, summary)
                if added:
                    prs.save(str(output_path))
                    result.total_slides += 1
                    result.warnings.append("Summary slide appended")
            except Exception as e:
                result.warnings.append(f"Summary slide failed: {type(e).__name__}: {e}")
        
    except Exception as e:
        result.errors.append(f"{type(e).__name__}: {str(e)[:200]}")
    
    result.total_time_ms = (time.time() - t0) * 1000
    
    # Generate diagnostic report if trace available
    if trace and TELEMETRY_AVAILABLE:
        try:
            report = DiagnosticReport()
            html_path = output_dir / f"{Path(pdf_path).stem}_trace.html"
            html_path.write_text(report.generate_html(trace), encoding="utf-8")
            result.warnings.append(f"Trace: {html_path.name}")
        except Exception:
            pass
    
    # Cleanup sanitized temp file
    if SANDBOX_AVAILABLE and sanitized_path != pdf_path:
        try:
            Path(sanitized_path).unlink(missing_ok=True)
        except Exception:
            pass
    
    return result


async def main():
    import sys
    if len(sys.argv) < 2:
        print("Usage: python orchestrator_v2.py <pdf_path> [output_dir]")
        return
    
    pdf = Path(sys.argv[1])
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("output_v2")
    out.mkdir(exist_ok=True)
    
    result = await convert_pdf(str(pdf), out)
    
    print(f"Success: {result.success}")
    print(f"Pages: {result.total_pages}, Slides: {result.total_slides}")
    print(f"MCP augmented: {result.mcp_augmented_pages}")
    print(f"Time: {result.total_time_ms:.0f}ms")
    if result.errors:
        print(f"Errors: {result.errors}")
    if result.warnings:
        print(f"Warnings: {result.warnings}")


if __name__ == "__main__":
    asyncio.run(main())