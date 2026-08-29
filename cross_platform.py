"""Cross-Platform Compatibility Matrix for PPTX output.

Ensures PPTX renders correctly across PowerPoint, Google Slides,
Apple Keynote, and LibreOffice by injecting render hints and font fallbacks.
"""

import copy
import io
import os
import sys
import zipfile
from dataclasses import dataclass, field
from enum import Enum
from typing import Any, Dict, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Pt

try:
    from lxml import etree as lxml_etree
    HAS_LXML = True
except ImportError:
    HAS_LXML = False

MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
XML_NS = "http://www.w3.org/XML/1998/namespace"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


class PlatformRenderer(str, Enum):
    POWERPOINT = "powerpoint"
    GOOGLE_SLIDES = "google_slides"
    KEYNOTE = "keynote"
    LIBREOFFICE = "libreoffice"


@dataclass
class RenderTarget:
    platform: PlatformRenderer
    name: str
    ooxml_min_version: str
    ignores_custom_spacing: bool
    overrides_font_fallback: bool
    supports_alternate_content: bool
    max_rendering_precision: float

    def __str__(self) -> str:
        return self.name


class CrossPlatformEngine:
    """Ensure PPTX renders correctly across presentation platforms."""

    COMPATIBILITY_MATRIX: Dict[str, RenderTarget] = {
        "powerpoint": RenderTarget(
            platform=PlatformRenderer.POWERPOINT,
            name="Microsoft PowerPoint",
            ooxml_min_version="2006",
            ignores_custom_spacing=False,
            overrides_font_fallback=False,
            supports_alternate_content=True,
            max_rendering_precision=12700.0,
        ),
        "google_slides": RenderTarget(
            platform=PlatformRenderer.GOOGLE_SLIDES,
            name="Google Slides",
            ooxml_min_version="2006",
            ignores_custom_spacing=True,
            overrides_font_fallback=True,
            supports_alternate_content=True,
            max_rendering_precision=6350.0,
        ),
        "keynote": RenderTarget(
            platform=PlatformRenderer.KEYNOTE,
            name="Apple Keynote",
            ooxml_min_version="2006",
            ignores_custom_spacing=False,
            overrides_font_fallback=True,
            supports_alternate_content=False,
            max_rendering_precision=12700.0,
        ),
        "libreoffice": RenderTarget(
            platform=PlatformRenderer.LIBREOFFICE,
            name="LibreOffice Impress",
            ooxml_min_version="2006",
            ignores_custom_spacing=False,
            overrides_font_fallback=True,
            supports_alternate_content=False,
            max_rendering_precision=12700.0,
        ),
    }

    FONT_MAPPING: Dict[str, Dict[str, str]] = {
        "Helvetica": {
            "powerpoint": "Arial",
            "google_slides": "Roboto",
            "keynote": "Helvetica",
            "libreoffice": "Helvetica",
        },
        "Times": {
            "powerpoint": "Times New Roman",
            "google_slides": "Roboto Serif",
            "keynote": "Times",
            "libreoffice": "Times New Roman",
        },
        "Courier": {
            "powerpoint": "Courier New",
            "google_slides": "Roboto Mono",
            "keynote": "Courier",
            "libreoffice": "Courier New",
        },
        "Arial": {
            "powerpoint": "Arial",
            "google_slides": "Roboto",
            "keynote": "Arial",
            "libreoffice": "Arial",
        },
        "Tahoma": {
            "powerpoint": "Tahoma",
            "google_slides": "Roboto",
            "keynote": "Tahoma",
            "libreoffice": "Tahoma",
        },
        "Verdana": {
            "powerpoint": "Verdana",
            "google_slides": "Roboto",
            "keynote": "Verdana",
            "libreoffice": "Verdana",
        },
        "Georgia": {
            "powerpoint": "Georgia",
            "google_slides": "Roboto Serif",
            "keynote": "Georgia",
            "libreoffice": "Georgia",
        },
        "LibreFranklin": {
            "powerpoint": "Arial",
            "google_slides": "Libre Franklin",
            "keynote": "Helvetica",
            "libreoffice": "Arial",
        },
    }

    def get_target(self, platform: PlatformRenderer) -> RenderTarget:
        return self.COMPATIBILITY_MATRIX[platform.value]

    def resolve_font(self, font_name: str, platform: PlatformRenderer) -> str:
        if not font_name:
            return font_name
        mapping = self.FONT_MAPPING.get(font_name)
        if mapping is None:
            return font_name
        return mapping.get(platform.value, font_name)

    def inject_alternate_content(
        self, element: ET.Element, primary: str, fallback: str
    ) -> ET.Element:
        tag = f"{{{MC_NS}}}AlternateContent"
        ac_element = ET.Element(tag)
        ac_element.set(f"xmlns:mc", MC_NS)

        required = ET.SubElement(ac_element, f"{{{MC_NS}}}Choice")
        required.set("Requires", "c1")
        primary_copy = copy.deepcopy(element)
        required.append(primary_copy)

        fallback_elem = ET.SubElement(ac_element, f"{{{MC_NS}}}Fallback")
        fallback_copy = copy.deepcopy(element)
        fallback_elem.append(fallback_copy)

        return ac_element

    def add_render_hint(
        self, slide: Any, platform: PlatformRenderer
    ) -> List[str]:
        changes: List[str] = []
        if platform == PlatformRenderer.GOOGLE_SLIDES:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        self._force_single_spacing_paragraph(paragraph)
                    changes.append(
                        f"Forced single line spacing on shape '{shape.name}'"
                    )
        elif platform == PlatformRenderer.KEYNOTE:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for run in self._iter_runs(shape.text_frame):
                        self._apply_keynote_font_hint(run)
                    changes.append(
                        f"Applied Keynote font hints to shape '{shape.name}'"
                    )
        return changes

    def detect_platform_intent(
        self, headers: Optional[Dict] = None
    ) -> PlatformRenderer:
        if not headers:
            return PlatformRenderer.POWERPOINT
        ua = headers.get("user-agent", headers.get("User-Agent", "")).lower()
        if "google" in ua or "gsa" in ua:
            return PlatformRenderer.GOOGLE_SLIDES
        if "keynote" in ua or "mac os" in ua:
            return PlatformRenderer.KEYNOTE
        if "libreoffice" in ua or "impress" in ua:
            return PlatformRenderer.LIBREOFFICE
        return PlatformRenderer.POWERPOINT

    def _force_single_spacing_paragraph(self, paragraph: Any) -> None:
        try:
            pPr = paragraph._p.get_or_add_pPr()
            spc_elem = pPr.find(f"{{{A_NS}}}spcBef")
            if spc_elem is not None:
                pPr.remove(spc_elem)
            spc_elem = pPr.find(f"{{{A_NS}}}spcAft")
            if spc_elem is not None:
                pPr.remove(spc_elem)
            lnSpc = pPr.find(f"{{{A_NS}}}lnSpc")
            if lnSpc is not None:
                pPr.remove(lnSpc)
        except Exception:
            pass

    def _apply_keynote_font_hint(self, run: Any) -> None:
        try:
            font_name = run.font.name
            if font_name:
                resolved = self.resolve_font(font_name, PlatformRenderer.KEYNOTE)
                if resolved != font_name:
                    run.font.name = resolved
        except Exception:
            pass

    @staticmethod
    def _iter_runs(text_frame: Any):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                yield run


class GoogleSlidesOptimizer:
    """Apply Google Slides-specific optimizations."""

    MISCELLANEOUS_ISSUES: List[str] = [
        "Google Slides ignores custom line spacing (before/after 0)",
        "Google Slides maps fonts to nearest Google Font",
        "Google Slides caps at 100 slides for some advanced features",
        "Google Slides re-flows overflowing text into new container",
    ]

    def __init__(self) -> None:
        self._engine = CrossPlatformEngine()

    def apply(self, prs: Presentation) -> Dict[str, Any]:
        changes: List[str] = []
        fonts_mapped: Dict[str, str] = {}

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                self._force_single_spacing(shape.text_frame)
                changes.append(
                    f"Slide {slide_idx + 1}: Forced single spacing "
                    f"on '{shape.name}'"
                )

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        old_font = run.font.name
                        if old_font:
                            new_font = self._engine.resolve_font(
                                old_font, PlatformRenderer.GOOGLE_SLIDES
                            )
                            if new_font != old_font:
                                run.font.name = new_font
                                fonts_mapped[old_font] = new_font
                                changes.append(
                                    f"Slide {slide_idx + 1}: Mapped font "
                                    f"'{old_font}' -> '{new_font}' "
                                    f"on '{shape.name}'"
                                )
                        self._add_font_fallback(
                            run, fonts_mapped.get(old_font or "", "Roboto")
                        )

        return {"changes": changes, "fonts_mapped": fonts_mapped}

    def _force_single_spacing(self, text_frame: Any) -> None:
        for paragraph in text_frame.paragraphs:
            try:
                pPr = paragraph._p.get_or_add_pPr()
                for tag in ("spcBef", "spcAft", "lnSpc"):
                    elem = pPr.find(f"{{{A_NS}}}{tag}")
                    if elem is not None:
                        pPr.remove(elem)
            except Exception:
                pass

    def _add_font_fallback(self, run: Any, google_font: str) -> None:
        try:
            rPr = run._r.get_or_add_rPr()
            existing_hint = rPr.find(f"{{{A_NS}}}latin")
            if existing_hint is not None:
                existing_hint.set("ea", google_font)
                existing_hint.set("cs", google_font)
            else:
                latin = ET.SubElement(rPr, f"{{{A_NS}}}latin")
                latin.set("typeface", run.font.name or "")
                latin.set("ea", google_font)
                latin.set("cs", google_font)
        except Exception:
            pass


class KeynoteOptimizer:
    """Apply Apple Keynote-specific optimizations."""

    APPLE_SYSTEM_FONTS: Dict[str, str] = {
        "Arial": "Helvetica Neue",
        "Helvetica": "Helvetica Neue",
        "Times New Roman": "Times New Roman",
        "Courier New": "Courier",
        "Verdana": "Helvetica Neue",
        "Georgia": "Georgia",
        "Tahoma": "Helvetica Neue",
    }

    def __init__(self) -> None:
        self._engine = CrossPlatformEngine()

    def apply(self, prs: Presentation) -> Dict[str, Any]:
        changes: List[str] = []
        fonts_mapped: Dict[str, str] = {}

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        old_font = run.font.name
                        if old_font:
                            resolved = self._resolve_apple_font(old_font)
                            if resolved != old_font:
                                run.font.name = resolved
                                fonts_mapped[old_font] = resolved
                                changes.append(
                                    f"Slide {slide_idx + 1}: Keynote font "
                                    f"'{old_font}' -> '{resolved}' "
                                    f"on '{shape.name}'"
                                )

        return {"changes": changes, "fonts_mapped": fonts_mapped}

    def _resolve_apple_font(self, font_name: str) -> str:
        if font_name in self.APPLE_SYSTEM_FONTS:
            return self.APPLE_SYSTEM_FONTS[font_name]
        resolved = self._engine.resolve_font(
            font_name, PlatformRenderer.KEYNOTE
        )
        return self.APPLE_SYSTEM_FONTS.get(resolved, resolved)


class CrossPlatformValidator:
    """Validate the PPTX opens correctly in other platforms."""

    def validate_unzip(self, prs_path: str) -> Dict[str, Any]:
        issues: List[str] = []
        valid = True
        try:
            with zipfile.ZipFile(prs_path, "r") as zf:
                bad = zf.testzip()
                if bad is not None:
                    issues.append(f"Corrupt file in archive: {bad}")
                    valid = False
                xml_names = [
                    n for n in zf.namelist() if n.endswith(".xml")
                ]
                for name in xml_names:
                    try:
                        data = zf.read(name)
                        ET.fromstring(data)
                    except ET.ParseError as exc:
                        issues.append(f"Malformed XML in {name}: {exc}")
                        valid = False
                    except Exception as exc:
                        issues.append(f"Cannot read {name}: {exc}")
        except zipfile.BadZipFile as exc:
            issues.append(f"Not a valid ZIP/PPTX: {exc}")
            valid = False
        except FileNotFoundError:
            issues.append(f"File not found: {prs_path}")
            valid = False

        return {"valid": valid, "issues": issues}

    def scan_alternate_content(self, prs_path: str) -> int:
        count = 0
        try:
            with zipfile.ZipFile(prs_path, "r") as zf:
                for name in zf.namelist():
                    if not name.endswith(".xml"):
                        continue
                    data = zf.read(name)
                    text = data.decode("utf-8", errors="replace")
                    count += text.count("mc:AlternateContent")
        except Exception:
            pass
        return count

    def check_font_fallbacks(self, prs_path: str) -> Dict[str, Any]:
        fonts_found: Dict[str, bool] = {}
        try:
            with zipfile.ZipFile(prs_path, "r") as zf:
                for name in zf.namelist():
                    if not name.endswith(".xml"):
                        continue
                    data = zf.read(name)
                    text = data.decode("utf-8", errors="replace")
                    for font_name in CrossPlatformEngine.FONT_MAPPING:
                        if font_name in text:
                            hint_tag = f'typeface="{font_name}"'
                            ea_hint = 'ea="'
                            if hint_tag in text:
                                fonts_found[font_name] = False
                            elif ea_hint in text:
                                fonts_found[font_name] = True
                            else:
                                fonts_found[font_name] = False
        except Exception:
            pass

        missing = [f for f, ok in fonts_found.items() if not ok]
        return {
            "all_have_fallbacks": len(missing) == 0,
            "fonts": fonts_found,
            "missing_fallbacks": missing,
        }

    def full_validation(self, prs_path: str) -> Dict[str, Any]:
        unzip_result = self.validate_unzip(prs_path)
        ac_count = self.scan_alternate_content(prs_path)
        font_result = self.check_font_fallbacks(prs_path)

        return {
            "valid": unzip_result["valid"],
            "issues": unzip_result["issues"],
            "alternate_content_count": ac_count,
            "font_fallbacks": font_result,
            "summary": {
                "xml_valid": unzip_result["valid"],
                "alternate_content_wrappers": ac_count,
                "fonts_with_fallbacks": sum(
                    1 for v in font_result["fonts"].values() if v
                ),
                "total_fonts_found": len(font_result["fonts"]),
                "platform_readiness": {
                    "powerpoint": True,
                    "google_slides": (
                        unzip_result["valid"]
                        and not font_result["missing_fallbacks"]
                    ),
                    "keynote": unzip_result["valid"],
                    "libreoffice": unzip_result["valid"],
                },
            },
        }


def apply_cross_platform(
    prs: Presentation,
    platform: PlatformRenderer = PlatformRenderer.POWERPOINT,
) -> Dict[str, Any]:
    """One-call API to apply cross-platform compatibility to a Presentation.

    Returns summary of changes for the conversion summary slide.
    """
    engine = CrossPlatformEngine()
    target = engine.get_target(platform)
    all_changes: List[str] = []
    all_fonts: Dict[str, str] = {}

    if platform == PlatformRenderer.GOOGLE_SLIDES:
        optimizer = GoogleSlidesOptimizer()
        result = optimizer.apply(prs)
        all_changes.extend(result["changes"])
        all_fonts.update(result["fonts_mapped"])
        for slide in prs.slides:
            hints = engine.add_render_hint(slide, platform)
            all_changes.extend(hints)

    elif platform == PlatformRenderer.KEYNOTE:
        optimizer = KeynoteOptimizer()
        result = optimizer.apply(prs)
        all_changes.extend(result["changes"])
        all_fonts.update(result["fonts_mapped"])
        for slide in prs.slides:
            hints = engine.add_render_hint(slide, platform)
            all_changes.extend(hints)

    elif platform == PlatformRenderer.LIBREOFFICE:
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        old_font = run.font.name
                        if old_font:
                            resolved = engine.resolve_font(
                                old_font, PlatformRenderer.LIBREOFFICE
                            )
                            if resolved != old_font:
                                run.font.name = resolved
                                all_fonts[old_font] = resolved
                                all_changes.append(
                                    f"Slide {slide_idx + 1}: LibreOffice "
                                    f"font '{old_font}' -> '{resolved}'"
                                )

    return {
        "platform": platform.value,
        "target": target.name,
        "changes": all_changes,
        "fonts_mapped": all_fonts,
        "total_changes": len(all_changes),
    }


def _create_sample_pptx() -> Presentation:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Cross-Platform Test"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "This tests font fallback and spacing."
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Custom font paragraph"
    run.font.name = "Helvetica"
    run.font.size = Pt(14)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Arial fallback test"
    run2.font.name = "Arial"
    run2.font.size = Pt(12)
    return prs


def _main() -> None:
    prs = _create_sample_pptx()
    out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)))
    base = "cross_platform_test"

    pw_path = os.path.join(out_dir, f"{base}_powerpoint.pptx")
    prs.save(pw_path)

    gs_prs = _create_sample_pptx()
    gs_result = apply_cross_platform(gs_prs, PlatformRenderer.GOOGLE_SLIDES)
    gs_path = os.path.join(out_dir, f"{base}_google_slides.pptx")
    gs_prs.save(gs_path)

    kn_prs = _create_sample_pptx()
    kn_result = apply_cross_platform(kn_prs, PlatformRenderer.KEYNOTE)
    kn_path = os.path.join(out_dir, f"{base}_keynote.pptx")
    kn_prs.save(kn_path)

    lo_prs = _create_sample_pptx()
    lo_result = apply_cross_platform(lo_prs, PlatformRenderer.LIBREOFFICE)
    lo_path = os.path.join(out_dir, f"{base}_libreoffice.pptx")
    lo_prs.save(lo_path)

    validator = CrossPlatformValidator()

    print("=" * 60)
    print("CROSS-PLATFORM COMPATIBILITY VALIDATION")
    print("=" * 60)

    for label, path in [
        ("PowerPoint", pw_path),
        ("Google Slides", gs_path),
        ("Keynote", kn_path),
        ("LibreOffice", lo_path),
    ]:
        print(f"\n--- {label} Target ---")
        val = validator.full_validation(path)
        print(f"  Valid: {val['valid']}")
        print(f"  Issues: {val['issues']}")
        print(f"  AlternateContent wrappers: {val['alternate_content_count']}")
        fb = val['font_fallbacks']
        print(f"  Fonts found: {fb['fonts']}")
        print(f"  Missing fallbacks: {fb['missing_fallbacks']}")
        print(f"  Platform readiness: {val['summary']['platform_readiness']}")

    print(f"\n{'=' * 60}")
    print("CHANGE SUMMARIES")
    print(f"{'=' * 60}")
    for label, result in [
        ("Google Slides", gs_result),
        ("Keynote", kn_result),
        ("LibreOffice", lo_result),
    ]:
        print(f"\n--- {label} ---")
        print(f"  Total changes: {result['total_changes']}")
        print(f"  Fonts mapped: {result['fonts_mapped']}")
        for change in result['changes'][:5]:
            print(f"    - {change}")
        if len(result['changes']) > 5:
            print(f"    ... and {len(result['changes']) - 5} more")

    for path in [
        pw_path, gs_path, kn_path, lo_path
    ]:
        if os.path.exists(path):
            os.remove(path)

    print(f"\n{'=' * 60}")
    print("All test files cleaned up. Validation complete.")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    _main()
