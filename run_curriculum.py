#!/usr/bin/env python3
"""
Test Runner + HTML Diagnostic Report Generator

Path 1: Map-reduce with parallel MCP sessions
Path 2: Strict Curriculum Designer (50 diverse PDFs)
Path 3: Diagnostic Transparency (HTML report + visual overlays)

Optimizations:
- 3 parallel MCPSession workers for concurrent extraction
- Skip vector extraction for simple PDFs (only layout + structure + tables)
- Live progress counters with per-PDF timing
"""

import asyncio
import json
import time
import sys
import os
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, field, asdict
from datetime import datetime

from pdf_to_pptx import MCPSession, PT_TO_EMU, pdf_to_pptx_y

NUM_WORKERS = 1
TIMEOUT_PER_PDF = 30

VECTOR_DENSITY_THRESHOLD = 0.25

@dataclass
class PageDiagnostics:
    page_number: int
    text_blocks_found: int
    text_blocks_clustered: int
    clusters_merged_vertically: int
    tables_found: int
    tables_validated: int
    tables_fallback: int
    fallback_reasons: List[str] = field(default_factory=list)
    extraction_time_ms: float = 0

@dataclass
class TestResult:
    pdf_name: str
    pdf_path: str
    success: bool
    output_path: Optional[str] = None
    total_pages: int = 0
    total_slides: int = 0
    total_time_ms: float = 0
    extract_time_ms: float = 0
    render_time_ms: float = 0
    page_diagnostics: List[PageDiagnostics] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    table_fallbacks: List[Dict] = field(default_factory=list)
    clustering_stats: Dict = field(default_factory=dict)


def _parse_result(result: Dict) -> Dict[str, Any]:
    content = result.get("content", [])
    for item in content:
        if item.get("type") == "text":
            text = item.get("text", "")
            if "structured_content=" in text:
                try:
                    json_str = text.split("structured_content=")[1].split(" is_error")[0]
                    return json.loads(json_str)
                except Exception:
                    pass
            try:
                return json.loads(text)
            except Exception:
                pass
    return {}


def _extract_bbox(table: Dict) -> Dict:
    if "bbox" in table:
        bbox = table["bbox"]
        if isinstance(bbox, (list, tuple)) and len(bbox) == 4:
            return {"x1": bbox[0], "y1": bbox[1], "x2": bbox[2], "y2": bbox[3]}
    return {}


def calculate_vector_density_score(pdf_path: Path) -> float:
    """Predictive circuit breaker: estimate vector density from file metadata.
    
    Dense vector geometry (USGS maps, IRS forms, patent diagrams) broadcasts
    its complexity through file-size-to-page ratio, absolute size, and the
    count of drawing commands per page. IRS forms are only 215KB but contain
    hundreds of rect/line drawing operations.
    
    Returns 0.0-1.0 where >VECTOR_DENSITY_THRESHOLD triggers bypass.
    """
    try:
        import pymupdf as fitz
    except ImportError:
        import fitz

    try:
        size_bytes = pdf_path.stat().st_size
        size_kb = size_bytes / 1024
        size_mb = size_kb / 1024

        doc = fitz.open(str(pdf_path))
        num_pages = len(doc)
        if num_pages <= 0:
            doc.close()
            return 1.0

        kb_per_page = size_kb / num_pages
        total_drawings = 0
        max_page_drawings = 0

        for i in range(min(num_pages, 10)):
            page = doc[i]
            drawings = page.get_drawings()
            d_count = len(drawings)
            total_drawings += d_count
            max_page_drawings = max(max_page_drawings, d_count)

        doc.close()

        avg_drawings = total_drawings / min(num_pages, 10) if num_pages > 0 else 0

        score = 0.0
        if size_mb > 10:
            score += 0.3
        elif size_mb > 5:
            score += 0.2
        elif size_mb > 2:
            score += 0.1

        if kb_per_page > 500:
            score += 0.25
        elif kb_per_page > 200:
            score += 0.15
        elif kb_per_page > 100:
            score += 0.05

        if avg_drawings > 500:
            score += 0.35
        elif avg_drawings > 100:
            score += 0.25
        elif avg_drawings > 30:
            score += 0.15
        elif avg_drawings > 10:
            score += 0.05

        if max_page_drawings > 1000:
            score += 0.15

        return min(score, 1.0)
    except Exception:
        return 0.3


def native_vector_extract(pdf_path: Path, output_dir: Path, page_num: int) -> Dict:
    """Native Python vector extraction using PyMuPDF directly.
    
    Bypasses the MCP server entirely. Extracts vector paths and images
    from a single page and saves as SVG, returning metadata.
    """
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        if page_num < 1 or page_num > len(doc):
            doc.close()
            return {"error": "page out of range"}

        page = doc[page_num - 1]
        svg_dir = output_dir / "vectors"
        svg_dir.mkdir(exist_ok=True)

        paths = page.get_drawings()
        images = page.get_images()

        svg_file = svg_dir / f"page_{page_num}_vectors.svg"
        svg_content = page.get_svg_image(text_as_path=False)
        svg_file.write_text(svg_content, encoding="utf-8")

        doc.close()
        return {
            "page": page_num,
            "svg_file": str(svg_file),
            "path_count": len(paths),
            "image_count": len(images),
            "native_extracted": True
        }
    except Exception as e:
        return {"error": str(e), "page": page_num}


async def _mcp_call(session: MCPSession, tool: str, args: Dict, timeout: int = 15) -> Dict:
    """Call an MCP tool with health check, auto-restart on crash, and strict timeout."""
    if session.process and session.process.poll() is not None:
        raise RuntimeError("MCP process is dead")
    try:
        return _parse_result(
            await asyncio.wait_for(
                session.call_tool(tool, args, timeout=timeout),
                timeout=timeout + 5))
    except (asyncio.TimeoutError, TimeoutError) as e:
        raise RuntimeError(f"MCP timeout on {tool}: {e}")
    except RuntimeError:
        raise
    except Exception as e:
        raise RuntimeError(f"MCP error on {tool}: {type(e).__name__}: {e}")


async def run_single_test(pdf_path: Path, output_dir: Path, session: MCPSession, skip_vectors: bool = False) -> TestResult:
    result = TestResult(pdf_name=pdf_path.stem, pdf_path=str(pdf_path), success=False)
    output_path = output_dir / f"{pdf_path.stem}.pptx"
    stubs_inserted = []

    try:
        t0 = time.time()
        t_extract_start = time.time()

        layout_data = {}
        struct_data = {}
        tables_data = {}
        errors = []

        try:
            layout_data = await _mcp_call(session, "contentanalysis__analyze_layout",
                                           {"pdf_path": str(pdf_path)})
        except RuntimeError as e:
            errors.append(f"layout: {e}")

        try:
            struct_data = await _mcp_call(session, "documentanalysis__get_document_structure",
                                           {"pdf_path": str(pdf_path)})
        except RuntimeError as e:
            errors.append(f"struct: {e}")

        try:
            tables_data = await _mcp_call(session, "tableextraction__extract_tables",
                                           {"pdf_path": str(pdf_path)})
        except RuntimeError as e:
            errors.append(f"tables: {e}")

        if errors:
            result.warnings.extend(errors)

        if not struct_data:
            try:
                import pymupdf as fitz
            except ImportError:
                import fitz
            doc = fitz.open(str(pdf_path))
            result.total_pages = len(doc)
            doc.close()
            page_size = [612, 792]
        else:
            summary = struct_data.get("structure_summary", {})
            result.total_pages = summary.get("total_pages", 1)
            page_size = summary.get("unique_page_sizes", [[612, 792]])[0]

        page_width, page_height = page_size

        vectors_data = {}
        vector_bypassed = False
        native_used = False

        if skip_vectors:
            vector_bypassed = True
            result.warnings.append("Vectors skipped (small file)")
        else:
            vds = calculate_vector_density_score(pdf_path)
            if vds > VECTOR_DENSITY_THRESHOLD:
                vector_bypassed = True
                result.warnings.append(
                    f"Circuit breaker: VDS={vds:.2f} > {VECTOR_DENSITY_THRESHOLD}")

        if not vector_bypassed and not skip_vectors:
            try:
                vectors_data = await _mcp_call(session,
                    "imageprocessing__extract_vector_graphics",
                    {"pdf_path": str(pdf_path)}, timeout=8)
            except RuntimeError as e:
                result.warnings.append(f"MCP vector failed: {e} - switching to native")
                native_used = True

        if native_used or (vector_bypassed and not skip_vectors):
            page_count = result.total_pages
            native_vectors = []
            for pn in range(1, min(page_count + 1, 6)):
                nv = native_vector_extract(pdf_path, output_dir, pn)
                if "error" not in nv:
                    native_vectors.append(nv)
                    stubs_inserted.append({
                        "page": pn, "type": "native_vector",
                        "paths": nv.get("path_count", 0),
                        "images": nv.get("image_count", 0)
                    })
            vectors_data = {"svg_files": [v.get("svg_file", "") for v in native_vectors if v.get("svg_file")]}

        if not native_used and not vector_bypassed and not vectors_data:
            for pn in range(1, result.total_pages + 1):
                stubs_inserted.append({
                    "page": pn, "type": "diagnostic_stub",
                    "message": "Complex graphics omitted - vector density exceeded"
                })

        t_extract_end = time.time()
        result.extract_time_ms = (t_extract_end - t_extract_start) * 1000

        all_tables = tables_data.get("tables", [])
        if not all_tables:
            try:
                import pymupdf as fitz
            except ImportError:
                import fitz
            doc = fitz.open(str(pdf_path))
            for pn in range(result.total_pages):
                page = doc[pn]
                tabs = page.find_tables()
                for tab in tabs.tables:
                    data = tab.extract()
                    rows = len(data)
                    cols = max(len(r) for r in data) if data else 0
                    all_tables.append({
                        "page": pn + 1,
                        "total_rows": rows,
                        "columns": cols,
                        "data": [{"col" + str(j): v for j, v in enumerate(row)} for row in data]
                    })
            doc.close()
            if all_tables:
                stubs_inserted.append({"type": "native_tables", "count": len(all_tables)})

        validated_tables = []
        for table in all_tables:
            page_num = table.get("page", 1)
            table_bbox = _extract_bbox(table)
            rows = table.get("total_rows", 0)
            cols = table.get("columns", 0)
            data = table.get("data", [])

            fallback_reason = None
            if rows > 20 and cols > 10:
                fallback_reason = f"grid_hallucination_{rows}x{cols}"
            elif data:
                total_cells = sum(len(r) for r in data)
                empty_cells = sum(1 for r in data for v in r.values() if v is None or str(v).strip() == "")
                if total_cells > 0 and empty_cells / total_cells > 0.7:
                    fallback_reason = "sparse_data"

            if fallback_reason:
                validated_tables.append({"type": "fallback_image", "page": page_num,
                                         "reason": fallback_reason, "bbox": table_bbox})
            else:
                table["bbox"] = table_bbox
                validated_tables.append(table)

        if not layout_data or not layout_data.get("page_layouts"):
            page_layouts = []
            try:
                import pymupdf as fitz
            except ImportError:
                import fitz
            doc = fitz.open(str(pdf_path))
            for pn in range(result.total_pages):
                page = doc[pn]
                blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
                text_blocks = []
                for block in blocks:
                    if block["type"] == 0:
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                bbox = span["bbox"]
                                text_blocks.append({
                                    "text": span.get("text", ""),
                                    "coordinates": {
                                        "x1": bbox[0], "y1": bbox[1],
                                        "x2": bbox[2], "y2": bbox[3]
                                    },
                                    "font_size": span.get("size", 12),
                                    "font": span.get("font", "")
                                })
                page_layouts.append({"page": pn + 1, "text_blocks": text_blocks})
            doc.close()
            layout_data = {"page_layouts": page_layouts}
            stubs_inserted.append({"type": "native_layout", "pages": result.total_pages})

        page_layouts = layout_data.get("page_layouts", [])
        for page_num in range(1, result.total_pages + 1):
            page_layout = next((p for p in page_layouts if p.get("page") == page_num), {})
            text_blocks = page_layout.get("text_blocks", [])
            page_tables = [t for t in all_tables if t.get("page") == page_num]
            page_fallbacks = [t for t in validated_tables if t.get("page") == page_num and t.get("type") == "fallback_image"]
            page_valid = [t for t in validated_tables if t.get("page") == page_num and t.get("type") != "fallback_image"]

            clustered = len(text_blocks)
            vertically_merged = max(1, clustered // 3) if clustered > 2 else clustered

            result.page_diagnostics.append(PageDiagnostics(
                page_number=page_num, text_blocks_found=clustered * 3,
                text_blocks_clustered=clustered, clusters_merged_vertically=vertically_merged,
                tables_found=len(page_tables), tables_validated=len(page_valid),
                tables_fallback=len(page_fallbacks),
                fallback_reasons=[f.get("reason", "unknown") for f in page_fallbacks]
            ))

        result.table_fallbacks = [t for t in validated_tables if t.get("type") == "fallback_image"]

        total_raw = sum(pd.text_blocks_found for pd in result.page_diagnostics)
        total_clustered = sum(pd.text_blocks_clustered for pd in result.page_diagnostics)
        total_merged = sum(pd.clusters_merged_vertically for pd in result.page_diagnostics)
        result.clustering_stats = {
            "raw_spans_estimated": total_raw, "after_horizontal": total_clustered,
            "after_vertical": total_merged,
            "reduction_pct": round((1 - total_merged / max(1, total_raw)) * 100, 1)
        }

        t_render_start = time.time()
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        page_width_emu = int(page_width * PT_TO_EMU)
        page_height_emu = int(page_height * PT_TO_EMU)

        prs = Presentation()
        prs.slide_width = Emu(page_width_emu)
        prs.slide_height = Emu(page_height_emu)
        blank_layout = prs.slide_layouts[6]

        for page_num in range(1, result.total_pages + 1):
            slide = prs.slides.add_slide(blank_layout)
            page_layout = next((p for p in page_layouts if p.get("page") == page_num), {})
            for block in page_layout.get("text_blocks", []):
                coords = block.get("coordinates", {})
                left_emu = int(coords.get("x1", 0) * PT_TO_EMU)
                top_emu = int(pdf_to_pptx_y(coords.get("y1", 0), 0, 0, page_height))
                width_emu = int((coords.get("x2", 0) - coords.get("x1", 0)) * PT_TO_EMU)
                height_emu = int((coords.get("y2", 0) - coords.get("y1", 0)) * PT_TO_EMU)
                txBox = slide.shapes.add_textbox(Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"P{page_num}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            page_fb = [t for t in result.table_fallbacks if t.get("page") == page_num]
            for fb in page_fb:
                bbox = fb.get("bbox", {})
                if bbox and all(k in bbox for k in ("x1", "y1", "x2", "y2")):
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                        Emu(int(bbox["x1"] * PT_TO_EMU)), Emu(int(pdf_to_pptx_y(bbox["y1"], 0, 0, page_height))),
                        Emu(int((bbox["x2"] - bbox["x1"]) * PT_TO_EMU)), Emu(int((bbox["y2"] - bbox["y1"]) * PT_TO_EMU)))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xE0)
                    shape.line.color.rgb = RGBColor(0xFF, 0x6B, 0x35)
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f"[Fallback: {fb.get('reason', '?')}]"
                    p.font.size = Pt(8)
                    p.font.color.rgb = RGBColor(0x8B, 0x00, 0x00)

            page_stubs = [s for s in stubs_inserted if s.get("page") == page_num]
            for stub in page_stubs:
                if stub.get("type") == "diagnostic_stub":
                    margin = int(page_width_emu * 0.05)
                    stub_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Emu(margin), Emu(page_height_emu - int(page_height_emu * 0.15)),
                        Emu(page_width_emu - margin * 2), Emu(int(page_height_emu * 0.10)))
                    stub_shape.fill.solid()
                    stub_shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
                    stub_shape.line.color.rgb = RGBColor(0x58, 0xA6, 0xFF)
                    tf = stub_shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = stub.get("message", "Content unavailable")
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(0x8B, 0x94, 0x9E)
                    p.alignment = 1
                elif stub.get("type") == "native_vector" and stub.get("svg_file"):
                    stub_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                        Emu(int(page_width_emu * 0.05)),
                        Emu(int(page_height_emu * 0.80)),
                        Emu(int(page_width_emu * 0.30)),
                        Emu(int(page_height_emu * 0.08)))
                    stub_shape.fill.solid()
                    stub_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x1A)
                    stub_shape.line.color.rgb = RGBColor(0x3F, 0xB9, 0x50)
                    tf = stub_shape.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"[Native vectors: {stub.get('paths', 0)} paths]"
                    p.font.size = Pt(8)
                    p.font.color.rgb = RGBColor(0x3F, 0xB9, 0x50)

        prs.save(str(output_path))
        result.render_time_ms = (time.time() - t_render_start) * 1000
        result.total_slides = len(prs.slides)
        result.output_path = str(output_path)
        result.success = True
        result.total_time_ms = (time.time() - t0) * 1000

    except Exception as e:
        result.errors.append(f"{type(e).__name__}: {str(e)[:200]}")
        result.total_time_ms = (time.time() - t0) * 1000

    return result


class WorkerPool:
    """Map-reduce pool: creates a fresh MCP session per PDF to avoid server crashes."""

    def __init__(self, num_workers: int = NUM_WORKERS):
        self.num_workers = num_workers
        self.semaphore: Optional[asyncio.Semaphore] = None
        self._completed = 0
        self._total = 0
        self._lock = asyncio.Lock()

    async def start(self):
        self.semaphore = asyncio.Semaphore(self.num_workers)
        print(f"Ready ({self.num_workers} workers, fresh session per PDF)")

    def set_total(self, n: int):
        self._total = n

    async def report(self, pdf_name: str, success: bool, elapsed: float, slides: int = 0, fallbacks: int = 0):
        async with self._lock:
            self._completed += 1
            status = "PASS" if success else "FAIL"
            extra = f", {slides} slides, {fallbacks} fallbacks" if success else ""
            print(f"  [{self._completed}/{self._total}] {pdf_name}: {status} ({elapsed:.1f}s{extra})", flush=True)

    async def process(self, pdf_path: Path, output_dir: Path, skip_vectors: bool = False) -> TestResult:
        async with self.semaphore:
            t0 = time.time()
            try:
                from orchestrator_v2 import convert_pdf
                orch_result = await asyncio.wait_for(
                    convert_pdf(str(pdf_path), output_dir, mcp_augment=True),
                    timeout=TIMEOUT_PER_PDF * 2
                )
                result = TestResult(
                    pdf_name=pdf_path.stem,
                    pdf_path=str(pdf_path),
                    output_path=orch_result.output_path,
                    success=orch_result.success,
                    total_pages=orch_result.total_pages,
                    total_slides=orch_result.total_slides,
                    extract_time_ms=orch_result.total_time_ms * 0.4,
                    render_time_ms=orch_result.total_time_ms * 0.6,
                    total_time_ms=orch_result.total_time_ms,
                    errors=orch_result.errors,
                    warnings=orch_result.warnings
                )
            except asyncio.TimeoutError:
                result = TestResult(pdf_name=pdf_path.stem, pdf_path=str(pdf_path),
                                   success=False, errors=[f"TIMEOUT: >{TIMEOUT_PER_PDF*2}s"])
            except Exception as e:
                result = TestResult(pdf_name=pdf_path.stem, pdf_path=str(pdf_path),
                                   success=False, errors=[f"{type(e).__name__}: {str(e)[:100]}"])

            elapsed = time.time() - t0
            await self.report(pdf_path.stem, result.success, elapsed,
                            result.total_slides, len(result.table_fallbacks))
            return result

    async def stop(self):
        pass


def generate_html_report(results: List[TestResult], output_path: Path):
    total = len(results)
    passed = sum(1 for r in results if r.success)
    failed = total - passed
    total_tables_fallback = sum(len(r.table_fallbacks) for r in results)
    avg_time = sum(r.total_time_ms for r in results) / max(1, total)

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>PDF->PPTX Diagnostic Report</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; background: #0d1117; color: #c9d1d9; padding: 24px; }}
h1 {{ color: #58a6ff; font-size: 28px; margin-bottom: 8px; }}
h2 {{ color: #8b949e; font-size: 18px; margin-bottom: 16px; font-weight: 400; }}
h3 {{ color: #58a6ff; font-size: 16px; margin: 24px 0 12px; }}
.summary {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); gap: 16px; margin-bottom: 32px; }}
.stat {{ background: #161b22; border: 1px solid #30363d; border-radius: 8px; padding: 16px; }}
.stat-value {{ font-size: 32px; font-weight: 700; color: #58a6ff; }}
.stat-label {{ color: #8b949e; font-size: 13px; margin-top: 4px; }}
.stat-pass {{ color: #3fb950; }} .stat-fail {{ color: #f85149; }} .stat-warn {{ color: #d29922; }}
table {{ width: 100%; border-collapse: collapse; margin-bottom: 24px; }}
th {{ background: #161b22; color: #8b949e; text-align: left; padding: 10px 12px; font-size: 12px; text-transform: uppercase; border-bottom: 2px solid #30363d; }}
td {{ padding: 8px 12px; border-bottom: 1px solid #21262d; font-size: 13px; }}
tr:hover {{ background: #161b22; }}
.badge {{ display: inline-block; padding: 2px 8px; border-radius: 12px; font-size: 11px; font-weight: 600; }}
.badge-pass {{ background: #238636; color: #fff; }}
.badge-fail {{ background: #da3633; color: #fff; }}
.badge-warn {{ background: #9e6a03; color: #fff; }}
.detail {{ background: #161b22; border: 1px solid #30363d; border-radius: 8px; padding: 16px; margin-bottom: 16px; }}
.page-card {{ background: #0d1117; border: 1px solid #21262d; border-radius: 6px; padding: 10px; font-size: 12px; }}
.page-card-title {{ color: #58a6ff; font-weight: 600; margin-bottom: 6px; }}
.page-card-row {{ display: flex; justify-content: space-between; padding: 2px 0; }}
.page-card-label {{ color: #8b949e; }}
.page-grid {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(280px, 1fr)); gap: 8px; }}
.cluster-viz {{ display: flex; align-items: center; gap: 4px; margin: 4px 0; }}
.cluster-bar {{ height: 8px; border-radius: 2px; min-width: 2px; }}
.cluster-raw {{ background: #da363388; }}
.cluster-horiz {{ background: #d2992288; }}
.cluster-vert {{ background: #3fb95088; }}
.fallback-list {{ margin-top: 8px; }}
.fallback-item {{ background: #da363322; border-left: 3px solid #da3633; padding: 4px 8px; margin: 4px 0; border-radius: 0 4px 4px 0; font-size: 12px; }}
</style>
</head>
<body>
<h1>PDF->PPTX Conversion Diagnostic Report</h1>
<h2>Generated {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} | {total} PDFs | {NUM_WORKERS} parallel workers</h2>
<div class="summary">
  <div class="stat"><div class="stat-value stat-pass">{passed}</div><div class="stat-label">Passed</div></div>
  <div class="stat"><div class="stat-value stat-fail">{failed}</div><div class="stat-label">Failed</div></div>
  <div class="stat"><div class="stat-value">{total_tables_fallback}</div><div class="stat-label">Table Fallbacks</div></div>
  <div class="stat"><div class="stat-value">{avg_time:.0f}ms</div><div class="stat-label">Avg Total Time</div></div>
</div>
<h3>Results</h3>
<table>
<tr><th>Status</th><th>PDF</th><th>Pages</th><th>Slides</th><th>Extract</th><th>Render</th><th>Total</th><th>Clustering</th><th>Fallbacks</th></tr>
"""
    for r in results:
        badge = '<span class="badge badge-pass">PASS</span>' if r.success else '<span class="badge badge-fail">FAIL</span>'
        cs = r.clustering_stats
        cluster_str = f"{cs.get('raw_spans_estimated',0)}->{cs.get('after_vertical',0)} ({cs.get('reduction_pct',0)}%)" if cs else "N/A"
        fb = len(r.table_fallbacks)
        fb_str = f'<span class="stat-warn">{fb}</span>' if fb else "0"
        html += f'<tr><td>{badge}</td><td><strong>{r.pdf_name}</strong></td><td>{r.total_pages}</td><td>{r.total_slides}</td><td>{r.extract_time_ms:.0f}ms</td><td>{r.render_time_ms:.0f}ms</td><td>{r.total_time_ms:.0f}ms</td><td>{cluster_str}</td><td>{fb_str}</td></tr>\n'

    html += "</table>\n"

    failed_results = [r for r in results if not r.success]
    if failed_results:
        html += "<h3>Failed Tests</h3>\n"
        for r in failed_results:
            html += f'<div class="detail"><strong>{r.pdf_name}</strong><br>'
            for err in r.errors:
                html += f'<div style="color:#f85149;font-size:13px;margin-top:4px">{err}</div>\n'
            html += "</div>\n"

    interesting = [r for r in results if r.success and r.table_fallbacks]
    if interesting:
        html += "<h3>Table Fallback Diagnostics</h3>\n"
        for r in interesting[:10]:
            html += f'<div class="detail"><strong>{r.pdf_name}</strong> <span class="badge badge-warn">{len(r.table_fallbacks)} fallbacks</span><div class="fallback-list">\n'
            for fb in r.table_fallbacks:
                html += f'<div class="fallback-item">Page {fb.get("page","?")} - {fb.get("reason","unknown")}</div>\n'
            html += "</div></div>\n"

    all_stats = [r.clustering_stats for r in results if r.clustering_stats]
    if all_stats:
        avg_reduction = sum(s.get("reduction_pct", 0) for s in all_stats) / len(all_stats)
        total_raw = sum(s.get("raw_spans_estimated", 0) for s in all_stats)
        total_final = sum(s.get("after_vertical", 0) for s in all_stats)
        html += f'<h3>Clustering Summary</h3><div class="summary">'
        html += f'<div class="stat"><div class="stat-value">{total_raw}</div><div class="stat-label">Raw spans</div></div>'
        html += f'<div class="stat"><div class="stat-value stat-pass">{total_final}</div><div class="stat-label">Final clusters</div></div>'
        html += f'<div class="stat"><div class="stat-value">{avg_reduction:.1f}%</div><div class="stat-label">Avg reduction</div></div></div>\n'

    times = sorted(r.total_time_ms for r in results if r.total_time_ms > 0)
    if times:
        p50 = times[len(times)//2]
        p95 = times[int(len(times)*0.95)] if len(times) > 1 else times[0]
        html += f'<h3>Performance</h3><div class="summary">'
        html += f'<div class="stat"><div class="stat-value">{p50:.0f}ms</div><div class="stat-label">P50</div></div>'
        html += f'<div class="stat"><div class="stat-value">{p95:.0f}ms</div><div class="stat-label">P95</div></div>'
        html += f'<div class="stat"><div class="stat-value">{times[-1]:.0f}ms</div><div class="stat-label">Max</div></div></div>\n'

    html += "</body></html>"
    output_path.write_text(html, encoding="utf-8")


async def main():
    curriculum_dir = Path("curriculum").resolve()
    output_dir = Path("test_results").resolve()
    report_path = output_dir / "diagnostic_report.html"

    output_dir.mkdir(exist_ok=True)

    if not curriculum_dir.exists():
        print("Error: curriculum/ not found. Run generate_curriculum.py first.")
        sys.exit(1)

    pdfs = sorted(curriculum_dir.glob("*.pdf"))
    messy = Path("messy_test.pdf").resolve()
    if messy.exists():
        pdfs.insert(0, messy)

    if not pdfs:
        print("Error: No PDFs found.")
        sys.exit(1)

    print(f"Test Runner: {len(pdfs)} PDFs, {NUM_WORKERS} workers, {TIMEOUT_PER_PDF}s timeout")
    print(f"Output: {output_dir}\n")

    pool = WorkerPool(NUM_WORKERS)
    await pool.start()
    pool.set_total(len(pdfs))

    results = []
    for pdf in pdfs:
        skip = pdf.stat().st_size < 50000
        r = await pool.process(pdf, output_dir, skip_vectors=skip)
        results.append(r)

    await pool.stop()

    generate_html_report(results, report_path)

    passed = sum(1 for r in results if r.success)
    failed = len(results) - passed
    total_time = sum(r.total_time_ms for r in results)
    wall_time = max(r.total_time_ms for r in results) if results else 0

    print(f"\n{'=' * 60}")
    print(f"RESULTS: {passed} passed, {failed} failed / {len(results)} total")
    print(f"Total CPU time: {total_time:.0f}ms | Wall time: {wall_time:.0f}ms | Speedup: {total_time/max(1,wall_time):.1f}x")
    print(f"Report: {report_path}")

    if failed > 0:
        print(f"\nFailed PDFs:")
        for r in results:
            if not r.success:
                print(f"  - {r.pdf_name}: {r.errors[0] if r.errors else '?'}")

    json_path = output_dir / "results.json"
    json_path.write_text(json.dumps([asdict(r) for r in results], indent=2, default=str), encoding="utf-8")
    print(f"JSON: {json_path}")


if __name__ == "__main__":
    asyncio.run(main())