#!/usr/bin/env python3
"""
theme_mapper.py — Semantic Theme Mapper for PDF-to-PPTX Converter

Classifies extracted text clusters into semantic roles and snaps them to
native PPTX theme elements. Replaces hardcoded absolute positioning and
RGB colors with theme-aware rendering.

Standalone module — no imports from orchestrator_v2.py or pdf_to_pptx.py.
"""

import re
import math
import zipfile
import xml.etree.ElementTree as ET
from enum import Enum
from typing import Dict, List, Optional, Tuple, Any
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

PT_TO_EMU = 12700

OOXML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

_LIST_PATTERNS = re.compile(
    r"^\s*[-\u2022\u25CF\u25CB\u25A0\u25B6]\s+"
    r"|^\s*\d+[.)]\s+"
    r"|^\s*[a-zA-Z][.)]\s+"
)


# ---------------------------------------------------------------------------
# SemanticRole Enum
# ---------------------------------------------------------------------------

class SemanticRole(Enum):
    TITLE = "title"
    SUBTITLE = "subtitle"
    BODY = "body"
    CAPTION = "caption"
    HEADER = "header"
    FOOTER = "footer"
    LIST_ITEM = "list_item"
    COLUMN_HEADER = "column_header"


# ---------------------------------------------------------------------------
# ThemeColors — reads & snaps to PPTX theme palette
# ---------------------------------------------------------------------------

class ThemeColors:
    """Extracts the color palette from a PPTX theme.xml and provides
    nearest-match snapping via Euclidean distance in RGB space."""

    _DEFAULT_THEME = {
        "dk1":      (0x00, 0x00, 0x00),
        "lt1":      (0xFF, 0xFF, 0xFF),
        "dk2":      (0x44, 0x54, 0x6A),
        "lt2":      (0xE7, 0xE6, 0xE6),
        "accent1":  (0x44, 0x72, 0xC4),
        "accent2":  (0xED, 0x7D, 0x31),
        "accent3":  (0xA5, 0xA5, 0xA5),
        "accent4":  (0xFF, 0xC0, 0x00),
        "accent5":  (0x5B, 0x9B, 0xD5),
        "accent6":  (0x70, 0xAD, 0x47),
        "hlink":    (0x05, 0x63, 0xC1),
        "folHlink": (0x95, 0x47, 0x91),
    }

    def __init__(self, prs: Presentation):
        self.colors: Dict[str, Tuple[int, int, int]] = {}
        self._load_from_presentation(prs)

    # -- public API ---------------------------------------------------------

    def snap_to_nearest(self, hex_color: str) -> Tuple[str, str]:
        """Return (theme_color_name, hex_value) closest to *hex_color*."""
        if not self.colors:
            return ("", hex_color)
        r, g, b = self._hex_to_rgb(hex_color)
        best_name = ""
        best_dist = float("inf")
        for name, (tr, tg, tb) in self.colors.items():
            dist = math.sqrt((r - tr) ** 2 + (g - tg) ** 2 + (b - tb) ** 2)
            if dist < best_dist:
                best_dist = dist
                best_name = name
        snapped = self.colors[best_name]
        return (best_name, self._rgb_to_hex(*snapped))

    # -- internals ----------------------------------------------------------

    def _load_from_presentation(self, prs: Presentation) -> None:
        """Try to read theme.xml from the PPTX ZIP; fall back to defaults."""
        try:
            prs_part = prs.part
            part_names = [
                "/ppt/theme/theme1.xml",
                "/ppt/theme/theme2.xml",
            ]
            theme_bytes: Optional[bytes] = None
            # python-pptx stores the package — walk its parts
            for rel in prs_part.rels.values():
                try:
                    pn = rel.target_ref
                    if pn and "theme" in pn and pn.endswith(".xml"):
                        theme_bytes = rel.target_part.blob
                        break
                except Exception:
                    continue
            if theme_bytes is None:
                self.colors = dict(self._DEFAULT_THEME)
                return
            self._parse_theme_xml(theme_bytes)
        except Exception:
            self.colors = dict(self._DEFAULT_THEME)

    def _parse_theme_xml(self, xml_bytes: bytes) -> None:
        """Parse <a:clrScheme> from theme XML bytes."""
        root = ET.fromstring(xml_bytes)
        clr_scheme = root.find(".//a:clrScheme", OOXML_NS)
        if clr_scheme is None:
            self.colors = dict(self._DEFAULT_THEME)
            return
        self.colors.clear()
        for child in clr_scheme:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            rgb_elem = child.find("a:srgbClr", OOXML_NS)
            if rgb_elem is not None:
                val = rgb_elem.get("val", "")
                if val:
                    self.colors[tag] = self._hex_to_rgb(val)
                    continue
            sys_elem = child.find("a:sysClr", OOXML_NS)
            if sys_elem is not None:
                val = sys_elem.get("lastClr", "000000")
                self.colors[tag] = self._hex_to_rgb(val)
        if not self.colors:
            self.colors = dict(self._DEFAULT_THEME)

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
        h = hex_color.lstrip("#")
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _rgb_to_hex(r: int, g: int, b: int) -> str:
        return f"#{r:02X}{g:02X}{b:02X}"


# ---------------------------------------------------------------------------
# SemanticClassifier — assigns roles + snaps colors
# ---------------------------------------------------------------------------

class SemanticClassifier:
    """Classifies clustered text blocks into semantic roles.

    Input:  list of dicts from cluster_spans (keys: text, left, top, width,
            height, font_size, font_name …).
    Output: same list with ``_role`` and ``_theme_snapped_color`` added.
    """

    def __init__(
        self,
        page_height: float,
        body_font_size: float,
        table_bboxes: Optional[List[Tuple[float, float, float, float]]] = None,
        theme_colors: Optional[ThemeColors] = None,
    ):
        self.page_height = page_height
        self.body_font_size = body_font_size
        self.table_bboxes = table_bboxes or []
        self.theme = theme_colors

    def classify(self, clusters: List[Dict]) -> List[Dict]:
        """Classify all clusters and annotate in-place, returning the list."""
        if not clusters:
            return clusters

        sorted_by_font = sorted(clusters, key=lambda c: c.get("font_size", 0), reverse=True)

        assigned: Dict[int, SemanticRole] = {}

        # 1. TITLE — largest font in top 20% of page
        for c in sorted_by_font:
            top = c.get("top", 0)
            if top <= self.page_height * 0.20:
                assigned[id(c)] = SemanticRole.TITLE
                c["_role"] = SemanticRole.TITLE.value
                break

        # 2. SUBTITLE — second-largest font in top 30%
        for c in sorted_by_font:
            if id(c) in assigned:
                continue
            top = c.get("top", 0)
            if top <= self.page_height * 0.30:
                assigned[id(c)] = SemanticRole.SUBTITLE
                c["_role"] = SemanticRole.SUBTITLE.value
                break

        # 3. COLUMN_HEADER — bold text above a table bbox
        for c in clusters:
            if id(c) in assigned:
                continue
            if self._is_bold_cluster(c) and self._is_above_table(c):
                assigned[id(c)] = SemanticRole.COLUMN_HEADER
                c["_role"] = SemanticRole.COLUMN_HEADER.value

        # 4. LIST_ITEM — bullet-like patterns
        for c in clusters:
            if id(c) in assigned:
                continue
            text = c.get("text", "")
            if _LIST_PATTERNS.match(text):
                assigned[id(c)] = SemanticRole.LIST_ITEM
                c["_role"] = SemanticRole.LIST_ITEM.value

        # 5. HEADER — text in top 8% with font > body
        for c in clusters:
            if id(c) in assigned:
                continue
            top = c.get("top", 0)
            fs = c.get("font_size", 0)
            if top <= self.page_height * 0.08 and fs > self.body_font_size:
                assigned[id(c)] = SemanticRole.HEADER
                c["_role"] = SemanticRole.HEADER.value

        # 6. FOOTER — small text in bottom 8%
        for c in clusters:
            if id(c) in assigned:
                continue
            bottom = c.get("top", 0) + c.get("height", 0)
            fs = c.get("font_size", 0)
            if bottom >= self.page_height * 0.92 and fs < self.body_font_size * 0.85:
                assigned[id(c)] = SemanticRole.FOOTER
                c["_role"] = SemanticRole.FOOTER.value

        # 7. CAPTION — small font (< body * 0.85) in bottom 15%
        for c in clusters:
            if id(c) in assigned:
                continue
            top = c.get("top", 0)
            fs = c.get("font_size", 0)
            if top >= self.page_height * 0.85 and fs < self.body_font_size * 0.85:
                assigned[id(c)] = SemanticRole.CAPTION
                c["_role"] = SemanticRole.CAPTION.value

        # 8. Everything else → BODY
        for c in clusters:
            if id(c) not in assigned:
                c["_role"] = SemanticRole.BODY.value

        # Snap colors
        for c in clusters:
            c["_theme_snapped_color"] = self._resolve_color(c)

        return clusters

    # -- helpers ------------------------------------------------------------

    def _is_bold_cluster(self, cluster: Dict) -> bool:
        font_name = cluster.get("font_name", "")
        return "bold" in font_name.lower()

    def _is_above_table(self, cluster: Dict) -> bool:
        c_bottom = cluster.get("top", 0) + cluster.get("height", 0)
        c_left = cluster.get("left", 0)
        c_right = c_left + cluster.get("width", 0)
        for tb in self.table_bboxes:
            t_left, t_top = tb[0], tb[1]
            t_right = tb[2] if len(tb) > 2 else tb[0] + tb[2]
            if c_bottom <= t_top and c_left < t_right and c_right > t_left:
                return True
        return False

    def _resolve_color(self, cluster: Dict) -> str:
        """Return the theme-snapped hex color for a cluster."""
        raw = cluster.get("color", "")
        if isinstance(raw, int):
            # PyMuPDF integer color (grayscale or BGR packed)
            if raw == 0:
                hex_c = "#000000"
            else:
                hex_c = f"#{raw & 0xFF:02X}{(raw >> 8) & 0xFF:02X}{(raw >> 16) & 0xFF:02X}"
        elif isinstance(raw, str) and raw.startswith("#"):
            hex_c = raw
        else:
            hex_c = "#000000"
        if self.theme:
            _, snapped = self.theme.snap_to_nearest(hex_c)
            return snapped
        return hex_c


# ---------------------------------------------------------------------------
# Convenience function
# ---------------------------------------------------------------------------

def classify_and_enhance(
    clusters: List[Dict],
    page_height: float,
    body_font_size: float,
    table_bboxes: Optional[List[Tuple[float, float, float, float]]] = None,
    theme_colors: Optional[ThemeColors] = None,
) -> List[Dict]:
    """Classify clusters into semantic roles and snap their colors to the
    nearest theme color.  Returns the same list, mutated with ``_role``
    and ``_theme_snapped_color`` keys on each dict."""
    classifier = SemanticClassifier(
        page_height=page_height,
        body_font_size=body_font_size,
        table_bboxes=table_bboxes,
        theme_colors=theme_colors,
    )
    return classifier.classify(clusters)


# ---------------------------------------------------------------------------
# ThemeAwareRenderer
# ---------------------------------------------------------------------------

class ThemeAwareRenderer:
    """Renders classified clusters into a PPTX using native layouts and
    theme colors instead of hardcoded positioning / RGB."""

    def __init__(self, prs: Presentation, theme: ThemeColors):
        self.prs = prs
        self.theme = theme
        self._slide_width = prs.slide_width
        self._slide_height = prs.slide_height

    # -- public API ---------------------------------------------------------

    def render_page(
        self,
        clusters: List[Dict],
        page_width_pt: float,
        page_height_pt: float,
    ) -> None:
        """Add one slide to self.prs for a single page's clusters."""
        title = self._find_role(clusters, SemanticRole.TITLE)
        subtitle = self._find_role(clusters, SemanticRole.SUBTITLE)
        body_clusters = self._find_all_roles(
            clusters,
            [SemanticRole.BODY, SemanticRole.LIST_ITEM, SemanticRole.COLUMN_HEADER],
        )
        caption_clusters = self._find_all_roles(
            clusters,
            [SemanticRole.CAPTION, SemanticRole.FOOTER, SemanticRole.HEADER],
        )

        has_title = title is not None
        has_body = len(body_clusters) > 0

        layout = self._pick_layout(has_title, has_body)
        slide = self.prs.slides.add_slide(layout)

        if has_title:
            self._place_title(slide, title, layout)
            if subtitle is not None:
                self._place_subtitle(slide, subtitle, layout)

        if has_body:
            self._place_body(slide, body_clusters, layout, page_width_pt, page_height_pt)

        for cc in caption_clusters:
            self._place_floating(slide, cc, page_width_pt, page_height_pt)

    # -- layout selection ---------------------------------------------------

    def _pick_layout(self, has_title: bool, has_body: bool) -> Any:
        layouts = self.prs.slide_layouts
        if has_title and has_body:
            # Index 5 = Title + Content
            if len(layouts) > 5:
                return layouts[5]
        if has_title and not has_body:
            # Index 0 = Title Slide
            if len(layouts) > 0:
                return layouts[0]
        # Fallback: Blank
        if len(layouts) > 6:
            return layouts[6]
        return layouts[-1]

    # -- content placement --------------------------------------------------

    def _place_title(self, slide, cluster: Dict, layout) -> None:
        placeholder = self._find_placeholder(slide, layout, expect_title=True)
        if placeholder is not None:
            tf = placeholder.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = cluster.get("text", "")
            color_hex = cluster.get("_theme_snapped_color", "#000000")
            p.font.color.rgb = self._hex_to_rgb_color(color_hex)
            fs = cluster.get("font_size", 24)
            p.font.size = Pt(fs)
        else:
            self._place_floating(slide, cluster, 0, 0)

    def _place_subtitle(self, slide, cluster: Dict, layout) -> None:
        placeholder = self._find_placeholder(slide, layout, expect_subtitle=True)
        if placeholder is not None:
            tf = placeholder.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = cluster.get("text", "")
            color_hex = cluster.get("_theme_snapped_color", "#333333")
            p.font.color.rgb = self._hex_to_rgb_color(color_hex)
            fs = cluster.get("font_size", 18)
            p.font.size = Pt(fs)
        else:
            self._place_floating(slide, cluster, 0, 0)

    def _place_body(
        self,
        slide,
        clusters: List[Dict],
        layout,
        page_width_pt: float,
        page_height_pt: float,
    ) -> None:
        placeholder = self._find_body_placeholder(slide, layout)
        if placeholder is not None:
            tf = placeholder.text_frame
            tf.clear()
            for i, cluster in enumerate(clusters):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = cluster.get("text", "")
                color_hex = cluster.get("_theme_snapped_color", "#333333")
                p.font.color.rgb = self._hex_to_rgb_color(color_hex)
                fs = cluster.get("font_size", 12)
                p.font.size = Pt(fs)
                if cluster.get("_role") == SemanticRole.LIST_ITEM.value:
                    p.level = 0
        else:
            for cluster in clusters:
                self._place_floating(slide, cluster, page_width_pt, page_height_pt)

    def _place_floating(
        self, slide, cluster: Dict, page_width_pt: float, page_height_pt: float,
    ) -> None:
        """Fall back to absolute-positioned textbox (existing approach)."""
        bbox = (
            cluster.get("left", 0),
            cluster.get("top", 0),
            cluster.get("left", 0) + cluster.get("width", 0),
            cluster.get("top", 0) + cluster.get("height", 0),
        )
        left_emu = int(bbox[0] * PT_TO_EMU)
        top_emu = int((page_height_pt - bbox[1]) * PT_TO_EMU) if page_height_pt else int(bbox[1] * PT_TO_EMU)
        width_emu = int(max((bbox[2] - bbox[0]), 1) * PT_TO_EMU)
        height_emu = int(max((bbox[3] - bbox[1]), 1) * PT_TO_EMU)

        txBox = slide.shapes.add_textbox(Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cluster.get("text", "")
        color_hex = cluster.get("_theme_snapped_color", "#333333")
        p.font.color.rgb = self._hex_to_rgb_color(color_hex)
        fs = cluster.get("font_size", 12)
        p.font.size = Pt(fs)

    # -- placeholder helpers ------------------------------------------------

    def _find_placeholder(self, slide, layout, expect_title=False, expect_subtitle=False):
        from pptx.enum.shapes import PP_PLACEHOLDER
        for shape in slide.placeholders:
            ph = shape.placeholder_format
            if ph.idx < 0:
                continue
            if expect_title and ph.type == PP_PLACEHOLDER.TITLE:
                return shape
            if expect_subtitle and ph.type == PP_PLACEHOLDER.SUBTITLE:
                return shape
        return None

    def _find_body_placeholder(self, slide, layout):
        from pptx.enum.shapes import PP_PLACEHOLDER
        for shape in slide.placeholders:
            ph = shape.placeholder_format
            if ph.idx < 0:
                continue
            if ph.type == PP_PLACEHOLDER.OBJECT:
                return shape
            if ph.type == PP_PLACEHOLDER.BODY:
                return shape
            if ph.idx == 1:
                return shape
        return None

    # -- helpers ------------------------------------------------------------

    @staticmethod
    def _find_role(clusters: List[Dict], role: SemanticRole) -> Optional[Dict]:
        for c in clusters:
            if c.get("_role") == role.value:
                return c
        return None

    @staticmethod
    def _find_all_roles(clusters: List[Dict], roles: List[SemanticRole]) -> List[Dict]:
        role_vals = {r.value for r in roles}
        return [c for c in clusters if c.get("_role") in role_vals]

    @staticmethod
    def _hex_to_rgb_color(hex_str: str) -> RGBColor:
        h = hex_str.lstrip("#")
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# apply_theme_to_presentation — top-level convenience
# ---------------------------------------------------------------------------

def apply_theme_to_presentation(
    prs: Presentation,
    title_cluster: Optional[Dict],
    body_clusters: Optional[List[Dict]] = None,
    table_clusters: Optional[List[Dict]] = None,
) -> None:
    """Populate an existing Presentation with theme-aware content.

    *prs* should already have slide_width / slide_height set.
    This function adds one slide using the best-matching layout, applying
    theme colors and semantic roles.  Existing slides are left untouched.

    Parameters
    ----------
    prs : Presentation
        Target presentation (modified in-place).
    title_cluster : dict or None
        Cluster dict with at least ``text``, ``left``, ``top``, ``width``,
        ``height``, ``font_size`` keys.
    body_clusters : list of dict, optional
        Body / list-item clusters to render as content.
    table_clusters : list of dict, optional
        Table-formatted clusters (rendered as floating text for now).
    """
    theme = ThemeColors(prs)
    body_clusters = body_clusters or []
    table_clusters = table_clusters or []

    all_text = []
    if title_cluster:
        title_cluster.setdefault("_role", SemanticRole.TITLE.value)
        if "_theme_snapped_color" not in title_cluster:
            title_cluster["_theme_snapped_color"] = _quick_snap(theme, title_cluster)
        all_text.append(title_cluster)
    for bc in body_clusters:
        bc.setdefault("_role", SemanticRole.BODY.value)
        if "_theme_snapped_color" not in bc:
            bc["_theme_snapped_color"] = _quick_snap(theme, bc)
        all_text.append(bc)
    for tc in table_clusters:
        tc.setdefault("_role", SemanticRole.BODY.value)
        if "_theme_snapped_color" not in tc:
            tc["_theme_snapped_color"] = _quick_snap(theme, tc)
        all_text.append(tc)

    renderer = ThemeAwareRenderer(prs, theme)
    page_w = 612.0
    page_h = 792.0
    if hasattr(prs, "slide_width") and prs.slide_width:
        page_w = prs.slide_width / PT_TO_EMU
    if hasattr(prs, "slide_height") and prs.slide_height:
        page_h = prs.slide_height / PT_TO_EMU
    renderer.render_page(all_text, page_w, page_h)


def _quick_snap(theme: ThemeColors, cluster: Dict) -> str:
    raw = cluster.get("color", "")
    if isinstance(raw, int):
        if raw == 0:
            hex_c = "#000000"
        else:
            hex_c = f"#{raw & 0xFF:02X}{(raw >> 8) & 0xFF:02X}{(raw >> 16) & 0xFF:02X}"
    elif isinstance(raw, str) and raw.startswith("#"):
        hex_c = raw
    else:
        hex_c = "#000000"
    _, snapped = theme.snap_to_nearest(hex_c)
    return snapped
