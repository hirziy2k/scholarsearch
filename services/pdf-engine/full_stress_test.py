#!/usr/bin/env python3
"""
full_stress_test.py — Comprehensive Stress Test for OpenCode/OmniRoute/MCP Stack

Tests:
1. MCP Server Startup & Response Times
2. PDF Tools (extract_text, extract_images, extract_tables, etc.)
3. PowerPoint Tools (add_slide, add_animation, etc.)
4. Skill Execution
5. OmniRoute API Performance
6. Memory & Concurrency
7. End-to-End Pipeline
"""

import asyncio
import json
import time
import sys
import os
import statistics
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, field, asdict
from datetime import datetime
import subprocess

# ============================================================================
# CONFIGURATION
# ============================================================================

TEST_PDF = "messy_test.pdf"
TEST_OUTPUT_DIR = Path("test_results")
TEST_OUTPUT_DIR.mkdir(exist_ok=True)

def _load_env_key() -> str:
    """Read the OmniRoute API key from the env file (~/.omniroute/.env)."""
    env_file = Path.home() / ".omniroute" / ".env"
    try:
        for line in env_file.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if line.startswith("OMNIROUTE_API_KEY="):
                return line.split("=", 1)[1].strip().strip('"').strip("'")
    except (OSError, ValueError):
        return ""
    return ""


@dataclass
class TestResult:
    test_name: str
    category: str
    status: str  # "pass", "fail", "skip", "timeout"
    duration_ms: float = 0.0
    details: str = ""
    error: str = ""
    metadata: Dict = field(default_factory=dict)

@dataclass
class StressTestReport:
    start_time: datetime = field(default_factory=datetime.now)
    end_time: Optional[datetime] = None
    results: List[TestResult] = field(default_factory=list)
    
    @property
    def total(self) -> int:
        return len(self.results)
    
    @property
    def passed(self) -> int:
        return sum(1 for r in self.results if r.status == "pass")
    
    @property
    def failed(self) -> int:
        return sum(1 for r in self.results if r.status == "fail")
    
    @property
    def skipped(self) -> int:
        return sum(1 for r in self.results if r.status == "skip")
    
    @property
    def total_duration_ms(self) -> float:
        return sum(r.duration_ms for r in self.results)
    
    @property
    def avg_duration_ms(self) -> float:
        return self.total_duration_ms / max(1, self.total)
    
    def add(self, result: TestResult):
        self.results.append(result)
    
    def to_json(self) -> str:
        return json.dumps({
            "start_time": self.start_time.isoformat(),
            "end_time": self.end_time.isoformat() if self.end_time else None,
            "summary": {
                "total": self.total,
                "passed": self.passed,
                "failed": self.failed,
                "skipped": self.skipped,
                "pass_rate": f"{(self.passed/self.total*100):.1f}%" if self.total else "0%",
                "total_duration_ms": self.total_duration_ms,
                "avg_duration_ms": self.avg_duration_ms,
            },
            "results": [asdict(r) for r in self.results]
        }, indent=2)


report = StressTestReport()


# ============================================================================
# HELPERS
# ============================================================================

async def measure_time(coro):
    """Measure execution time of an async coroutine."""
    start = time.perf_counter()
    try:
        result = await coro
        duration = (time.perf_counter() - start) * 1000
        return result, duration, None
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000
        return None, duration, e


def run_command(cmd: List[str], timeout: int = 30) -> tuple:
    """Run a command and return (stdout, stderr, returncode, duration_ms)."""
    start = time.perf_counter()
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=str(Path(__file__).parent)
        )
        duration = (time.perf_counter() - start) * 1000
        return result.stdout, result.stderr, result.returncode, duration
    except subprocess.TimeoutExpired:
        duration = (time.perf_counter() - start) * 1000
        return "", "TIMEOUT", -1, duration
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000
        return "", str(e), -1, duration


# ============================================================================
# TEST CATEGORY 1: MCP Server Startup
# ============================================================================

async def test_mcp_startup():
    """Test MCP server startup times."""
    print("\n[1] MCP SERVER STARTUP TESTS")
    print("-" * 60)
    
    servers = [
        ("pdf-tools", ["uvx", "mcp-pdf"]),
        ("powerpoint", ["uvx", "--with", "mcp<2", "powerpoint-mcp"]),
    ]
    
    for name, cmd in servers:
        print(f"  Testing {name} startup...", end=" ", flush=True)
        
        # Test 1: Cold start
        stdout, stderr, rc, duration = run_command(cmd + ["--help"], timeout=15)
        
        if rc == 0 or "usage" in stdout.lower() or "usage" in stderr.lower():
            report.add(TestResult(
                test_name=f"mcp_startup_{name}",
                category="mcp_startup",
                status="pass",
                duration_ms=duration,
                details=f"Cold start in {duration:.0f}ms",
                metadata={"server": name, "cold_start_ms": duration}
            ))
            print(f"PASS ({duration:.0f}ms)")
        else:
            report.add(TestResult(
                test_name=f"mcp_startup_{name}",
                category="mcp_startup",
                status="fail",
                duration_ms=duration,
                error=stderr[:200] if stderr else "Unknown error",
                metadata={"server": name, "returncode": rc}
            ))
            print(f"FAIL (rc={rc})")


# ============================================================================
# TEST CATEGORY 2: PDF Tools
# ============================================================================

async def test_pdf_tools():
    """Test PDF tools with various operations."""
    print("\n[2] PDF TOOLS TESTS")
    print("-" * 60)
    
    from pdf_extractor import PyMuPDFExtractor
    
    if not Path(TEST_PDF).exists():
        print(f"  SKIP: {TEST_PDF} not found")
        report.add(TestResult(
            test_name="pdf_tools_skip",
            category="pdf_tools",
            status="skip",
            details=f"{TEST_PDF} not found"
        ))
        return
    
    # Test 2.1: PyMuPDF extraction speed
    print("  [2.1] PyMuPDF extraction...", end=" ", flush=True)
    start = time.perf_counter()
    try:
        with PyMuPDFExtractor(TEST_PDF) as extractor:
            doc = extractor.extract_all()
            duration = (time.perf_counter() - start) * 1000
            
            page_count = len(doc.pages)
            total_blocks = sum(len(p.blocks) for p in doc.pages)
            total_tables = sum(len(p.tables) for p in doc.pages)
            
            report.add(TestResult(
                test_name="pdf_pymupdf_extraction",
                category="pdf_tools",
                status="pass",
                duration_ms=duration,
                details=f"{page_count} pages, {total_blocks} blocks, {total_tables} tables in {duration:.0f}ms",
                metadata={
                    "pages": page_count,
                    "blocks": total_blocks,
                    "tables": total_tables,
                    "ms_per_page": duration / max(1, page_count)
                }
            ))
            print(f"PASS ({duration:.0f}ms, {page_count} pages)")
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000
        report.add(TestResult(
            test_name="pdf_pymupdf_extraction",
            category="pdf_tools",
            status="fail",
            duration_ms=duration,
            error=str(e)[:200]
        ))
        print(f"FAIL ({e})")
    
    # Test 2.2: Complexity analysis
    print("  [2.2] Complexity analysis...", end=" ", flush=True)
    start = time.perf_counter()
    try:
        with PyMuPDFExtractor(TEST_PDF) as extractor:
            doc = extractor.extract_all()
            complexities = []
            for i in range(min(5, len(doc.pages))):  # Test first 5 pages
                c = extractor.analyze_complexity(extractor.doc[i])
                complexities.append(c)
            duration = (time.perf_counter() - start) * 1000
            
            complex_pages = sum(1 for c in complexities if c.get("complex"))
            report.add(TestResult(
                test_name="pdf_complexity_analysis",
                category="pdf_tools",
                status="pass",
                duration_ms=duration,
                details=f"{len(complexities)} pages analyzed, {complex_pages} complex",
                metadata={"complex_pages": complex_pages, "total_analyzed": len(complexities)}
            ))
            print(f"PASS ({duration:.0f}ms, {complex_pages} complex)")
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000
        report.add(TestResult(
            test_name="pdf_complexity_analysis",
            category="pdf_tools",
            status="fail",
            duration_ms=duration,
            error=str(e)[:200]
        ))
        print(f"FAIL ({e})")
    
    # Test 2.3: MCP-PDF session startup
    print("  [2.3] MCP-PDF session startup...", end=" ", flush=True)
    try:
        from pdf_to_pptx import MCPSession
        session = MCPSession("uvx", ["mcp-pdf"])
        start = time.perf_counter()
        await session.start(timeout=5.0)
        duration = (time.perf_counter() - start) * 1000
        session.close()
        
        report.add(TestResult(
            test_name="mcp_pdf_session_startup",
            category="pdf_tools",
            status="pass",
            duration_ms=duration,
            details=f"MCP-PDF session started in {duration:.0f}ms"
        ))
        print(f"PASS ({duration:.0f}ms)")
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000 if 'start' in dir() else 0
        report.add(TestResult(
            test_name="mcp_pdf_session_startup",
            category="pdf_tools",
            status="fail",
            duration_ms=duration,
            error=str(e)[:200]
        ))
        print(f"FAIL ({e})")


# ============================================================================
# TEST CATEGORY 3: PowerPoint Tools
# ============================================================================

async def test_powerpoint_tools():
    """Test PowerPoint MCP tools."""
    print("\n[3] POWERPOINT TOOLS TESTS")
    print("-" * 60)
    
    # Test 3.1: PowerPoint MCP startup
    print("  [3.1] PowerPoint MCP startup...", end=" ", flush=True)
    start = time.perf_counter()
    stdout, stderr, rc, duration = run_command(
        ["uvx", "--with", "mcp<2", "powerpoint-mcp", "--help"],
        timeout=15
    )
    
    if rc == 0 or "usage" in stdout.lower() or "usage" in stderr.lower():
        report.add(TestResult(
            test_name="powerpoint_mcp_startup",
            category="powerpoint_tools",
            status="pass",
            duration_ms=duration,
            details=f"PowerPoint MCP started in {duration:.0f}ms"
        ))
        print(f"PASS ({duration:.0f}ms)")
    else:
        report.add(TestResult(
            test_name="powerpoint_mcp_startup",
            category="powerpoint_tools",
            status="fail",
            duration_ms=duration,
            error=stderr[:200]
        ))
        print(f"FAIL (rc={rc})")
    
    # Test 3.2: python-pptx library
    print("  [3.2] python-pptx library...", end=" ", flush=True)
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        
        start = time.perf_counter()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(100000), Emu(50000))
        txBox.text_frame.paragraphs[0].text = "Test"
        prs.save(str(TEST_OUTPUT_DIR / "test.pptx"))
        duration = (time.perf_counter() - start) * 1000
        
        report.add(TestResult(
            test_name="python_pptx_library",
            category="powerpoint_tools",
            status="pass",
            duration_ms=duration,
            details=f"Created test PPTX in {duration:.0f}ms"
        ))
        print(f"PASS ({duration:.0f}ms)")
        
        # Cleanup
        (TEST_OUTPUT_DIR / "test.pptx").unlink(missing_ok=True)
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000 if 'start' in dir() else 0
        report.add(TestResult(
            test_name="python_pptx_library",
            category="powerpoint_tools",
            status="fail",
            duration_ms=duration,
            error=str(e)[:200]
        ))
        print(f"FAIL ({e})")


# ============================================================================
# TEST CATEGORY 4: OmniRoute API
# ============================================================================

async def test_omniroute_api():
    """Test OmniRoute API performance."""
    print("\n[4] OMNIRoute API TESTS")
    print("-" * 60)
    
    API_KEY = os.environ.get("OMNIROUTE_API_KEY", "") or _load_env_key()
    if not API_KEY:
        print("  SKIP: OMNIROUTE_API_KEY not set (see ~/.omniroute/.env)")
        print("  " + "-" * 60)
        return
    BASE_URL = os.environ.get("OMNIROUTE_BASE_URL", "http://127.0.0.1:20128")
    
    try:
        import httpx
    except ImportError:
        print("  SKIP: httpx not installed")
        report.add(TestResult(
            test_name="omniroute_skip",
            category="omniroute_api",
            status="skip",
            details="httpx not installed"
        ))
        return
    
    # Test 4.1: Health check latency
    print("  [4.1] Health check latency...", end=" ", flush=True)
    async with httpx.AsyncClient(timeout=5.0) as client:
        latencies = []
        for _ in range(5):
            start = time.perf_counter()
            try:
                resp = await client.get(f"{BASE_URL}/v1/models", headers={"Authorization": f"Bearer {API_KEY}"})
                latencies.append((time.perf_counter() - start) * 1000)
            except:
                latencies.append(9999)
        
        avg_latency = statistics.mean(latencies)
        p95_latency = sorted(latencies)[int(len(latencies) * 0.95)]
        
        status = "pass" if avg_latency < 500 else "fail"
        report.add(TestResult(
            test_name="omniroute_health_latency",
            category="omniroute_api",
            status=status,
            duration_ms=avg_latency,
            details=f"Avg: {avg_latency:.0f}ms, P95: {p95_latency:.0f}ms",
            metadata={"avg_ms": avg_latency, "p95_ms": p95_latency, "min_ms": min(latencies), "max_ms": max(latencies)}
        ))
        print(f"{'PASS' if status == 'pass' else 'FAIL'} (avg={avg_latency:.0f}ms, p95={p95_latency:.0f}ms)")
    
    # Test 4.2: Model listing
    print("  [4.2] Model listing...", end=" ", flush=True)
    async with httpx.AsyncClient(timeout=10.0) as client:
        start = time.perf_counter()
        try:
            resp = await client.get(f"{BASE_URL}/v1/models", headers={"Authorization": f"Bearer {API_KEY}"})
            duration = (time.perf_counter() - start) * 1000
            data = resp.json()
            model_count = len(data.get("data", []))
            
            report.add(TestResult(
                test_name="omniroute_model_listing",
                category="omniroute_api",
                status="pass" if resp.status_code == 200 else "fail",
                duration_ms=duration,
                details=f"{model_count} models listed in {duration:.0f}ms",
                metadata={"model_count": model_count, "status_code": resp.status_code}
            ))
            print(f"PASS ({model_count} models, {duration:.0f}ms)")
        except Exception as e:
            duration = (time.perf_counter() - start) * 1000
            report.add(TestResult(
                test_name="omniroute_model_listing",
                category="omniroute_api",
                status="fail",
                duration_ms=duration,
                error=str(e)[:200]
            ))
            print(f"FAIL ({e})")
    
    # Test 4.3: Chat completion (quick test)
    print("  [4.3] Chat completion latency...", end=" ", flush=True)
    async with httpx.AsyncClient(timeout=30.0) as client:
        start = time.perf_counter()
        try:
            resp = await client.post(
                f"{BASE_URL}/v1/chat/completions",
                headers={"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"},
                json={
                    "model": "north-mini-code-free",
                    "messages": [{"role": "user", "content": "Say hi in 5 words"}],
                    "max_tokens": 20
                }
            )
            duration = (time.perf_counter() - start) * 1000
            
            if resp.status_code == 200:
                data = resp.json()
                tokens = data.get("usage", {}).get("total_tokens", 0)
                report.add(TestResult(
                    test_name="omniroute_chat_completion",
                    category="omniroute_api",
                    status="pass",
                    duration_ms=duration,
                    details=f"Completed in {duration:.0f}ms, {tokens} tokens",
                    metadata={"tokens": tokens, "status_code": resp.status_code}
                ))
                print(f"PASS ({duration:.0f}ms, {tokens} tokens)")
            else:
                report.add(TestResult(
                    test_name="omniroute_chat_completion",
                    category="omniroute_api",
                    status="fail",
                    duration_ms=duration,
                    error=f"HTTP {resp.status_code}: {resp.text[:100]}",
                    metadata={"status_code": resp.status_code}
                ))
                print(f"FAIL (HTTP {resp.status_code})")
        except Exception as e:
            duration = (time.perf_counter() - start) * 1000
            report.add(TestResult(
                test_name="omniroute_chat_completion",
                category="omniroute_api",
                status="fail",
                duration_ms=duration,
                error=str(e)[:200]
            ))
            print(f"FAIL ({e})")


# ============================================================================
# TEST CATEGORY 5: Skill Execution
# ============================================================================

async def test_skills():
    """Test skill loading and execution."""
    print("\n[5] SKILL EXECUTION TESTS")
    print("-" * 60)
    
    skill_dir = Path(".opencode/skills")
    
    if not skill_dir.exists():
        print("  SKIP: .opencode/skills not found")
        report.add(TestResult(
            test_name="skills_skip",
            category="skills",
            status="skip",
            details=".opencode/skills not found"
        ))
        return
    
    # Test 5.1: Skill discovery
    print("  [5.1] Skill discovery...", end=" ", flush=True)
    start = time.perf_counter()
    skills = list(skill_dir.glob("**/*.md"))
    duration = (time.perf_counter() - start) * 1000
    
    report.add(TestResult(
        test_name="skill_discovery",
        category="skills",
        status="pass",
        duration_ms=duration,
        details=f"Found {len(skills)} skills",
        metadata={"skill_count": len(skills), "skills": [s.stem for s in skills]}
    ))
    print(f"PASS ({len(skills)} skills, {duration:.0f}ms)")
    
    # Test 5.2: Skill parse time
    print("  [5.2] Skill parse time...", end=" ", flush=True)
    parse_times = []
    for skill in skills:
        start = time.perf_counter()
        content = skill.read_text(encoding="utf-8")
        # Parse frontmatter
        if content.startswith("---"):
            parts = content.split("---", 2)
            if len(parts) >= 3:
                frontmatter = parts[1]
                body = parts[2]
        parse_times.append((time.perf_counter() - start) * 1000)
    
    avg_parse = statistics.mean(parse_times) if parse_times else 0
    report.add(TestResult(
        test_name="skill_parse_time",
        category="skills",
        status="pass",
        duration_ms=avg_parse,
        details=f"Avg parse time: {avg_parse:.2f}ms per skill",
        metadata={"avg_parse_ms": avg_parse}
    ))
    print(f"PASS (avg={avg_parse:.2f}ms)")
    
    # Test 5.3: Skill content size
    print("  [5.3] Skill content size...", end=" ", flush=True)
    sizes = []
    for skill in skills:
        size = skill.stat().st_size
        sizes.append(size)
    
    avg_size = statistics.mean(sizes) if sizes else 0
    total_size = sum(sizes)
    report.add(TestResult(
        test_name="skill_content_size",
        category="skills",
        status="pass",
        duration_ms=0,
        details=f"Avg: {avg_size:.0f} bytes, Total: {total_size:.0f} bytes",
        metadata={"avg_bytes": avg_size, "total_bytes": total_size}
    ))
    print(f"PASS (avg={avg_size:.0f} bytes, total={total_size:.0f} bytes)")


# ============================================================================
# TEST CATEGORY 6: Concurrency & Memory
# ============================================================================

async def test_concurrency():
    """Test concurrent operations."""
    print("\n[6] CONCURRENCY TESTS")
    print("-" * 60)
    
    # Test 6.1: Concurrent MCP startups
    print("  [6.1] Concurrent MCP startups...", end=" ", flush=True)
    
    async def startup_mcp(name, cmd):
        start = time.perf_counter()
        try:
            proc = await asyncio.create_subprocess_exec(
                *cmd, "--help",
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=10)
            duration = (time.perf_counter() - start) * 1000
            return name, duration, None
        except Exception as e:
            duration = (time.perf_counter() - start) * 1000
            return name, duration, str(e)
    
    tasks = [
        startup_mcp("pdf-tools-1", ["uvx", "mcp-pdf"]),
        startup_mcp("pdf-tools-2", ["uvx", "mcp-pdf"]),
        startup_mcp("pptx-tools-1", ["uvx", "--with", "mcp<2", "powerpoint-mcp"]),
        startup_mcp("pptx-tools-2", ["uvx", "--with", "mcp<2", "powerpoint-mcp"]),
    ]
    
    results = await asyncio.gather(*tasks)
    durations = [r[1] for r in results]
    errors = [r for r in results if r[2]]
    
    avg_duration = statistics.mean(durations)
    status = "pass" if len(errors) < 2 else "fail"
    
    report.add(TestResult(
        test_name="concurrent_mcp_startup",
        category="concurrency",
        status=status,
        duration_ms=avg_duration,
        details=f"4 concurrent startups, avg: {avg_duration:.0f}ms, errors: {len(errors)}",
        metadata={"avg_ms": avg_duration, "errors": len(errors)}
    ))
    print(f"{'PASS' if status == 'pass' else 'FAIL'} (avg={avg_duration:.0f}ms, {len(errors)} errors)")
    
    # Test 6.2: Memory baseline
    print("  [6.2] Memory baseline...", end=" ", flush=True)
    try:
        import psutil
        process = psutil.Process()
        mem_mb = process.memory_info().rss / (1024 * 1024)
        report.add(TestResult(
            test_name="memory_baseline",
            category="concurrency",
            status="pass",
            duration_ms=0,
            details=f"Current process memory: {mem_mb:.1f} MB",
            metadata={"memory_mb": mem_mb}
        ))
        print(f"PASS ({mem_mb:.1f} MB)")
    except ImportError:
        report.add(TestResult(
            test_name="memory_baseline",
            category="concurrency",
            status="skip",
            details="psutil not installed"
        ))
        print("SKIP (psutil not installed)")


# ============================================================================
# TEST CATEGORY 7: End-to-End Pipeline
# ============================================================================

async def test_e2e_pipeline():
    """Test full PDF-to-PPTX pipeline."""
    print("\n[7] END-TO-END PIPELINE TESTS")
    print("-" * 60)
    
    if not Path(TEST_PDF).exists():
        print(f"  SKIP: {TEST_PDF} not found")
        report.add(TestResult(
            test_name="e2e_skip",
            category="e2e_pipeline",
            status="skip",
            details=f"{TEST_PDF} not found"
        ))
        return
    
    # Test 7.1: Full conversion
    print("  [7.1] Full PDF-to-PPTX conversion...", end=" ", flush=True)
    try:
        from orchestrator_v2 import convert_pdf
        
        start = time.perf_counter()
        result = await convert_pdf(TEST_PDF, TEST_OUTPUT_DIR, mcp_augment=True)
        duration = (time.perf_counter() - start) * 1000
        
        status = "pass" if result.success else "fail"
        report.add(TestResult(
            test_name="e2e_full_conversion",
            category="e2e_pipeline",
            status=status,
            duration_ms=duration,
            details=f"{result.total_pages} pages -> {result.total_slides} slides in {duration:.0f}ms",
            metadata={
                "pages": result.total_pages,
                "slides": result.total_slides,
                "success": result.success,
                "mcp_augmented": result.mcp_augmented_pages,
                "errors": len(result.errors),
                "warnings": len(result.warnings)
            }
        ))
        print(f"{'PASS' if status == 'pass' else 'FAIL'} ({result.total_pages} pages -> {result.total_slides} slides, {duration:.0f}ms)")
        
        if result.errors:
            for e in result.errors[:3]:
                print(f"    Error: {e[:100]}")
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000 if 'start' in dir() else 0
        report.add(TestResult(
            test_name="e2e_full_conversion",
            category="e2e_pipeline",
            status="fail",
            duration_ms=duration,
            error=str(e)[:300]
        ))
        print(f"FAIL ({e})")
    
    # Test 7.2: Conversion without MCP
    print("  [7.2] Conversion without MCP...", end=" ", flush=True)
    try:
        from orchestrator_v2 import convert_pdf
        
        start = time.perf_counter()
        result = await convert_pdf(TEST_PDF, TEST_OUTPUT_DIR, mcp_augment=False)
        duration = (time.perf_counter() - start) * 1000
        
        status = "pass" if result.success else "fail"
        report.add(TestResult(
            test_name="e2e_no_mcp_conversion",
            category="e2e_pipeline",
            status=status,
            duration_ms=duration,
            details=f"{result.total_pages} pages -> {result.total_slides} slides (no MCP) in {duration:.0f}ms",
            metadata={
                "pages": result.total_pages,
                "slides": result.total_slides,
                "success": result.success
            }
        ))
        print(f"{'PASS' if status == 'pass' else 'FAIL'} ({result.total_pages} pages -> {result.total_slides} slides, {duration:.0f}ms)")
    except Exception as e:
        duration = (time.perf_counter() - start) * 1000 if 'start' in dir() else 0
        report.add(TestResult(
            test_name="e2e_no_mcp_conversion",
            category="e2e_pipeline",
            status="fail",
            duration_ms=duration,
            error=str(e)[:300]
        ))
        print(f"FAIL ({e})")


# ============================================================================
# MAIN
# ============================================================================

async def main():
    print("=" * 70)
    print("  OPencode / OmniRoute / MCP / SKILL / PROJECT — FULL STRESS TEST")
    print("=" * 70)
    print(f"  Start: {report.start_time.strftime('%Y-%m-%d %H:%M:%S')}")
    print(f"  PDF: {TEST_PDF}")
    print(f"  Output: {TEST_OUTPUT_DIR}")
    print("=" * 70)
    
    # Run all test categories
    await test_mcp_startup()
    await test_pdf_tools()
    await test_powerpoint_tools()
    await test_omniroute_api()
    await test_skills()
    await test_concurrency()
    await test_e2e_pipeline()
    
    report.end_time = datetime.now()
    
    # Summary
    print("\n" + "=" * 70)
    print("  STRESS TEST SUMMARY")
    print("=" * 70)
    print(f"  Total Tests:    {report.total}")
    print(f"  Passed:         {report.passed}")
    print(f"  Failed:         {report.failed}")
    print(f"  Skipped:        {report.skipped}")
    print(f"  Pass Rate:      {(report.passed/report.total*100):.1f}%" if report.total else "  Pass Rate: 0%")
    print(f"  Total Duration: {report.total_duration_ms:.0f}ms")
    print(f"  Avg Duration:   {report.avg_duration_ms:.0f}ms")
    print("=" * 70)
    
    # Save report
    report_path = TEST_OUTPUT_DIR / "stress_test_report.json"
    report_path.write_text(report.to_json(), encoding="utf-8")
    print(f"\n  Full report saved to: {report_path}")
    
    # Save readable summary
    summary_path = TEST_OUTPUT_DIR / "stress_test_summary.txt"
    with open(summary_path, "w", encoding="utf-8") as f:
        f.write("STRESS TEST SUMMARY\n")
        f.write("=" * 70 + "\n")
        f.write(f"Date: {report.start_time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
        
        # Group by category
        categories = {}
        for r in report.results:
            cat = r.category
            if cat not in categories:
                categories[cat] = []
            categories[cat].append(r)
        
        for cat, tests in categories.items():
            f.write(f"\n{cat.upper()}\n")
            f.write("-" * 40 + "\n")
            for t in tests:
                status_icon = "PASS" if t.status == "pass" else "FAIL" if t.status == "fail" else "SKIP"
                f.write(f"  [{status_icon}] {t.test_name}: {t.details}\n")
                if t.error:
                    f.write(f"         Error: {t.error[:100]}\n")
        
        f.write("\n" + "=" * 70 + "\n")
        f.write(f"Pass Rate: {(report.passed/report.total*100):.1f}%\n" if report.total else "Pass Rate: 0%\n")
    
    print(f"  Summary saved to: {summary_path}")
    
    return report.failed == 0


if __name__ == "__main__":
    success = asyncio.run(main())
    sys.exit(0 if success else 1)
