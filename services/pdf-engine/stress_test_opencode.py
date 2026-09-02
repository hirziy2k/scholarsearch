#!/usr/bin/env python3
"""OpenCode Stack Stress Test — MCP, Dependencies, File I/O, Concurrency."""

import os
import sys
import time
import json
import hashlib
import tempfile
import threading
import concurrent.futures
import statistics
from pathlib import Path
from datetime import datetime

# ============================================================================
# TEST CONFIGURATION
# ============================================================================

RESULTS = []
TEMP_DIR = Path(tempfile.mkdtemp(prefix="opencode_stress_"))
TEST_PDF = Path(__file__).parent / "messy_test.pdf"

def log(msg, level="INFO"):
    print(f"[{level}] {msg}")

def record_test(name, duration, status, details=None):
    RESULTS.append({
        "name": name,
        "duration_ms": round(duration * 1000, 2),
        "status": status,
        "details": details or {},
        "timestamp": datetime.now().isoformat(),
    })
    icon = "PASS" if status == "pass" else "FAIL" if status == "fail" else "WARN"
    log(f"{icon} {name}: {round(duration * 1000, 2)}ms")

# ============================================================================
# TEST 1: Python Core Imports
# ============================================================================

def test_core_imports():
    """Test importing all core dependencies under load."""
    log("=== Test 1: Core Import Stress ===")
    modules = [
        "fitz", "pptx", "fastapi", "uvicorn", "httpx",
        "redis", "boto3", "psutil", "reportlab", "pydantic",
        "json", "hashlib", "threading", "concurrent.futures",
        "sqlite3", "pathlib", "tempfile",
    ]
    
    results = {}
    for mod in modules:
        start = time.perf_counter()
        try:
            __import__(mod)
            elapsed = time.perf_counter() - start
            results[mod] = {"status": "ok", "time_ms": round(elapsed * 1000, 3)}
        except Exception as e:
            results[mod] = {"status": "error", "error": str(e)}
    
    successes = sum(1 for v in results.values() if v["status"] == "ok")
    avg_time = statistics.mean([v["time_ms"] for v in results.values() if v["status"] == "ok"])
    
    record_test(
        "core_imports",
        0,
        "pass" if successes == len(modules) else "fail",
        {"successes": successes, "total": len(modules), "avg_import_ms": round(avg_time, 3)}
    )
    return results

# ============================================================================
# TEST 2: PyMuPDF (fitz) Stress
# ============================================================================

def test_pymupdf_stress():
    """Stress test PyMuPDF with repeated opens/closes."""
    log("=== Test 2: PyMuPDF Stress ===")
    import fitz
    
    if not TEST_PDF.exists():
        record_test("pymupdf_stress", 0, "warn", {"reason": "messy_test.pdf not found"})
        return
    
    iterations = 10
    timings = []
    
    for i in range(iterations):
        start = time.perf_counter()
        try:
            doc = fitz.open(str(TEST_PDF))
            page_count = len(doc)
            for page in doc:
                text = page.get_text()
                blocks = page.get_text("dict")["blocks"]
            doc.close()
            elapsed = time.perf_counter() - start
            timings.append(elapsed)
        except Exception as e:
            log(f"  Iteration {i} failed: {e}", "ERROR")
    
    if timings:
        record_test(
            "pymupdf_stress",
            statistics.mean(timings),
            "pass",
            {
                "iterations": iterations,
                "avg_ms": round(statistics.mean(timings) * 1000, 2),
                "min_ms": round(min(timings) * 1000, 2),
                "max_ms": round(max(timings) * 1000, 2),
                "stdev_ms": round(statistics.stdev(timings) * 1000, 2) if len(timings) > 1 else 0,
                "pages": page_count,
            }
        )

# ============================================================================
# TEST 3: python-pptx Stress
# ============================================================================

def test_pptx_stress():
    """Stress test python-pptx with repeated create/write cycles."""
    log("=== Test 3: python-pptx Stress ===")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    iterations = 20
    timings = []
    
    for i in range(iterations):
        start = time.perf_counter()
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            for s in range(5):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = f"Stress Test Slide {s+1}"
                for b in range(10):
                    from pptx.util import Emu
                    txBox = slide.shapes.add_textbox(
                        Inches(1), Inches(1 + b * 0.5), Inches(11), Inches(0.4)
                    )
                    tf = txBox.text_frame
                    tf.text = f"Bullet point {b+1} with some content to test text rendering"
            
            out_path = TEMP_DIR / f"stress_{i}.pptx"
            prs.save(str(out_path))
            elapsed = time.perf_counter() - start
            timings.append(elapsed)
        except Exception as e:
            log(f"  Iteration {i} failed: {e}", "ERROR")
    
    # Cleanup
    for f in TEMP_DIR.glob("stress_*.pptx"):
        f.unlink()
    
    if timings:
        record_test(
            "pptx_stress",
            statistics.mean(timings),
            "pass",
            {
                "iterations": iterations,
                "slides_per_file": 5,
                "avg_ms": round(statistics.mean(timings) * 1000, 2),
                "min_ms": round(min(timings) * 1000, 2),
                "max_ms": round(max(timings) * 1000, 2),
            }
        )

# ============================================================================
# TEST 4: Hashing / Crypto Stress
# ============================================================================

def test_hashing_stress():
    """Stress test SHA-256 hashing (used by entity_anonymizer, buffer_broker)."""
    log("=== Test 4: Hashing Stress ===")
    
    data_sizes = [1024, 10240, 102400, 1048576]  # 1KB to 1MB
    iterations = 100
    results = {}
    
    for size in data_sizes:
        payload = os.urandom(size)
        timings = []
        for _ in range(iterations):
            start = time.perf_counter()
            hashlib.sha256(payload).hexdigest()
            timings.append(time.perf_counter() - start)
        results[f"{size//1024}KB"] = {
            "avg_us": round(statistics.mean(timings) * 1_000_000, 2),
            "throughput_mb_s": round((size / statistics.mean(timings)) / (1024*1024), 1),
        }
    
    record_test(
        "hashing_stress",
        statistics.mean(timings),
        "pass",
        {"iterations_per_size": iterations, "results": results}
    )

# ============================================================================
# TEST 5: SQLite WAL Stress (slide_state pattern)
# ============================================================================

def test_sqlite_wal_stress():
    """Stress test SQLite WAL mode with concurrent reads/writes."""
    log("=== Test 5: SQLite WAL Stress ===")
    import sqlite3
    
    db_path = TEMP_DIR / "stress_test.sqlite"
    conn = sqlite3.connect(str(db_path))
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("""
        CREATE TABLE IF NOT EXISTS jobs (
            id INTEGER PRIMARY KEY,
            doc_id TEXT,
            slide_idx INTEGER,
            status TEXT DEFAULT 'pending',
            data TEXT
        )
    """)
    conn.commit()
    conn.close()
    
    iterations = 500
    errors = []
    
    def writer(thread_id):
        conn = sqlite3.connect(str(db_path))
        conn.execute("PRAGMA journal_mode=WAL")
        for i in range(iterations // 5):
            try:
                conn.execute(
                    "INSERT INTO jobs (doc_id, slide_idx, status, data) VALUES (?, ?, ?, ?)",
                    (f"doc_{thread_id}", i, "pending", json.dumps({"thread": thread_id}))
                )
                conn.commit()
            except Exception as e:
                errors.append(str(e))
        conn.close()
    
    def reader():
        conn = sqlite3.connect(str(db_path))
        conn.execute("PRAGMA journal_mode=WAL")
        for _ in range(iterations // 10):
            try:
                rows = conn.execute("SELECT COUNT(*) FROM jobs").fetchone()
                conn.execute("SELECT * FROM jobs WHERE status = 'pending' LIMIT 10").fetchall()
            except Exception as e:
                errors.append(str(e))
        conn.close()
    
    start = time.perf_counter()
    with concurrent.futures.ThreadPoolExecutor(max_workers=6) as pool:
        futures = [pool.submit(writer, i) for i in range(5)]
        futures += [pool.submit(reader) for _ in range(3)]
        concurrent.futures.wait(futures)
    elapsed = time.perf_counter() - start
    
    # Get final count
    conn = sqlite3.connect(str(db_path))
    count = conn.execute("SELECT COUNT(*) FROM jobs").fetchone()[0]
    conn.close()
    db_path.unlink()
    
    record_test(
        "sqlite_wal_stress",
        elapsed,
        "pass" if not errors else "warn",
        {
            "writers": 5,
            "readers": 3,
            "inserts_per_writer": iterations // 5,
            "total_rows": count,
            "errors": len(errors),
            "throughput_rows_s": round(count / elapsed),
        }
    )

# ============================================================================
# TEST 6: Concurrent File I/O
# ============================================================================

def test_concurrent_io():
    """Stress test concurrent file writes (simulating ephemeral storage)."""
    log("=== Test 6: Concurrent File I/O Stress ===")
    
    file_count = 50
    file_size = 10240  # 10KB each
    errors = []
    
    def write_file(idx):
        try:
            path = TEMP_DIR / f"io_test_{idx}.dat"
            data = os.urandom(file_size)
            path.write_bytes(data)
            # Read back and verify
            read_back = path.read_bytes()
            if hashlib.md5(data).hexdigest() != hashlib.md5(read_back).hexdigest():
                errors.append(f"File {idx} corrupted")
            path.unlink()
        except Exception as e:
            errors.append(str(e))
    
    start = time.perf_counter()
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as pool:
        list(pool.map(write_file, range(file_count)))
    elapsed = time.perf_counter() - start
    
    record_test(
        "concurrent_io",
        elapsed,
        "pass" if not errors else "fail",
        {
            "files": file_count,
            "size_per_file": file_size,
            "total_mb": round((file_count * file_size) / (1024*1024), 2),
            "throughput_mb_s": round(((file_count * file_size) / elapsed) / (1024*1024), 1),
            "errors": errors,
        }
    )

# ============================================================================
# TEST 7: FastAPI + httpx Proxy Simulation
# ============================================================================

def test_fastapi_httpx():
    """Stress test FastAPI + httpx (proxy pattern)."""
    log("=== Test 7: FastAPI + httpx Stress ===")
    from fastapi import FastAPI
    from fastapi.testclient import TestClient
    import httpx
    
    app = FastAPI()
    
    @app.get("/test")
    async def test_endpoint():
        return {"status": "ok", "data": "x" * 1000}
    
    @app.post("/chat/completions")
    async def chat_completions(body: dict):
        return {
            "id": "stress-test",
            "object": "chat.completion",
            "choices": [{"message": {"role": "assistant", "content": "test"}}],
            "usage": {"prompt_tokens": 10, "completion_tokens": 5},
        }
    
    client = TestClient(app)
    iterations = 200
    timings = []
    errors = 0
    
    for i in range(iterations):
        start = time.perf_counter()
        try:
            resp = client.get("/test")
            assert resp.status_code == 200
            resp = client.post("/chat/completions", json={"model": "test", "messages": []})
            assert resp.status_code == 200
            timings.append(time.perf_counter() - start)
        except Exception:
            errors += 1
    
    record_test(
        "fastapi_httpx",
        statistics.mean(timings) if timings else 0,
        "pass" if errors == 0 else "warn",
        {
            "iterations": iterations,
            "errors": errors,
            "avg_ms": round(statistics.mean(timings) * 1000, 2) if timings else 0,
            "p95_ms": round(sorted(timings)[int(len(timings)*0.95)] * 1000, 2) if timings else 0,
            "rps": round(iterations / sum(timings)) if timings else 0,
        }
    )

# ============================================================================
# TEST 8: Memory Pressure
# ============================================================================

def test_memory_pressure():
    """Test memory behavior under load."""
    log("=== Test 8: Memory Pressure ===")
    import psutil
    
    process = psutil.Process(os.getpid())
    mem_before = process.memory_info().rss / (1024*1024)
    
    # Allocate and release memory
    chunks = []
    for i in range(100):
        chunks.append(os.urandom(1024 * 1024))  # 1MB each
    
    mem_peak = process.memory_info().rss / (1024*1024)
    del chunks
    
    import gc
    gc.collect()
    mem_after = process.memory_info().rss / (1024*1024)
    
    record_test(
        "memory_pressure",
        0,
        "pass",
        {
            "mem_before_mb": round(mem_before, 1),
            "mem_peak_mb": round(mem_peak, 1),
            "mem_after_mb": round(mem_after, 1),
            "allocated_mb": 100,
            "reclaimed_mb": round(mem_peak - mem_after, 1),
        }
    )

# ============================================================================
# TEST 9: Thread Pool Saturation
# ============================================================================

def test_thread_saturation():
    """Test thread pool under heavy load."""
    log("=== Test 9: Thread Pool Saturation ===")
    
    work_items = 200
    errors = []
    
    def cpu_work(x):
        """Simulate some CPU work."""
        result = 0
        for i in range(1000):
            result += (x * i) % 1000
            result = result % 100000
        return result
    
    start = time.perf_counter()
    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as pool:
        futures = [pool.submit(cpu_work, i) for i in range(work_items)]
        results = [f.result() for f in concurrent.futures.as_completed(futures)]
    elapsed = time.perf_counter() - start
    
    record_test(
        "thread_saturation",
        elapsed,
        "pass",
        {
            "work_items": work_items,
            "workers": 8,
            "completed": len(results),
            "throughput_items_s": round(work_items / elapsed),
        }
    )

# ============================================================================
# MAIN
# ============================================================================

def main():
    log("=" * 60)
    log("OpenCode Stack Stress Test")
    log(f"Started: {datetime.now().isoformat()}")
    log("=" * 60)
    
    test_start = time.perf_counter()
    
    test_core_imports()
    test_pymupdf_stress()
    test_pptx_stress()
    test_hashing_stress()
    test_sqlite_wal_stress()
    test_concurrent_io()
    test_fastapi_httpx()
    test_memory_pressure()
    test_thread_saturation()
    
    total_elapsed = time.perf_counter() - test_start
    
    # Cleanup
    import shutil
    shutil.rmtree(TEMP_DIR, ignore_errors=True)
    
    # Summary
    log("=" * 60)
    log("STRESS TEST SUMMARY")
    log("=" * 60)
    
    passed = sum(1 for r in RESULTS if r["status"] == "pass")
    warned = sum(1 for r in RESULTS if r["status"] == "warn")
    failed = sum(1 for r in RESULTS if r["status"] == "fail")
    
    log(f"Total: {len(RESULTS)} | Pass: {passed} | Warn: {warned} | Fail: {failed}")
    log(f"Duration: {round(total_elapsed, 2)}s")
    
    for r in RESULTS:
        icon = "PASS" if r["status"] == "pass" else "WARN" if r["status"] == "warn" else "FAIL"
        log(f"  [{icon}] {r['name']}: {r['duration_ms']}ms")
    
    # Write report
    report_path = Path(__file__).parent / "test_results" / "opencode_stress_report.json"
    report_path.parent.mkdir(exist_ok=True)
    with open(report_path, "w") as f:
        json.dump({
            "timestamp": datetime.now().isoformat(),
            "total_duration_s": round(total_elapsed, 2),
            "summary": {"pass": passed, "warn": warned, "fail": failed},
            "tests": RESULTS,
        }, f, indent=2)
    log(f"\nReport: {report_path}")

if __name__ == "__main__":
    main()
