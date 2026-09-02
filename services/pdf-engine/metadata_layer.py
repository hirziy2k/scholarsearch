#!/usr/bin/env python3
"""
metadata_layer.py — PDF Metadata Extraction & PPTX Binding

Captures hyperlinks, bookmarks/TOC, page labels, and image alt-text
that PyMuPDF's get_text("dict") silently strips, then binds them onto
a python-pptx Presentation via MetadataBinder.
"""

import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass, field

import fitz
from pptx import Presentation


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class Hyperlink:
    page: int
    uri: str
    bbox: Tuple[float, float, float, float]
    text: str


@dataclass
class Bookmark:
    title: str
    level: int
    page: int


@dataclass
class PDFMetadata:
    hyperlinks: List[Hyperlink] = field(default_factory=list)
    bookmarks: List[Bookmark] = field(default_factory=list)
    page_labels: Dict[int, str] = field(default_factory=dict)
    alt_texts: Dict[str, str] = field(default_factory=dict)


# ---------------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------------

class MetadataExtractor:
    """Extracts hyperlinks, bookmarks, page labels, and alt-text from a PDF."""

    def __init__(self, pdf_path: str):
        self.pdf_path = pdf_path
        self._doc = fitz.open(pdf_path)

    def __enter__(self):
        return self

    def __exit__(self, *args):
        self._doc.close()

    # -- public API ---------------------------------------------------------

    def extract_all(self) -> PDFMetadata:
        return PDFMetadata(
            hyperlinks=self._extract_hyperlinks(),
            bookmarks=self._extract_toc(),
            page_labels=self._extract_page_labels(),
            alt_texts=self._extract_alt_texts(),
        )

    def get_hyperlinks_for_page(self, page_num: int) -> List[Hyperlink]:
        return [h for h in self._extract_hyperlinks() if h.page == page_num]

    def get_toc(self) -> List[Bookmark]:
        return self._extract_toc()

    # -- internals ----------------------------------------------------------

    def _extract_hyperlinks(self) -> List[Hyperlink]:
        results: List[Hyperlink] = []
        for page_idx in range(len(self._doc)):
            page = self._doc[page_idx]
            try:
                annots = list(page.annots()) if page.annots() else []
            except Exception:
                annots = []
            for annot in annots:
                if annot.type[0] != fitz.PDF_ANNOT_LINK:
                    continue
                uri = annot.info.get("uri", "")
                if not uri:
                    try:
                        uri = annot.uri or ""
                    except Exception:
                        uri = ""
                if not uri:
                    continue
                rect = annot.rect
                bbox = (rect.x0, rect.y0, rect.x1, rect.y1)
                text = self._text_in_rect(page, rect)
                results.append(Hyperlink(
                    page=page_idx + 1,
                    uri=uri,
                    bbox=bbox,
                    text=text,
                ))
        return results

    def _extract_toc(self) -> List[Bookmark]:
        try:
            toc = self._doc.get_toc()
        except Exception:
            toc = []
        return [
            Bookmark(title=entry[1], level=entry[0], page=entry[2])
            for entry in toc
        ]

    def _extract_page_labels(self) -> Dict[int, str]:
        try:
            labels = self._doc.get_page_labels()
        except Exception:
            labels = None
        if not labels:
            return {}
        result: Dict[int, str] = {}
        for idx in range(len(self._doc)):
            try:
                lbl = labels[idx] if idx < len(labels) else ""
            except Exception:
                lbl = ""
            if lbl:
                result[idx + 1] = lbl
        return result

    def _extract_alt_texts(self) -> Dict[str, str]:
        """Best-effort alt-text via image annotation metadata."""
        alt: Dict[str, str] = {}
        for page_idx in range(len(self._doc)):
            page = self._doc[page_idx]
            try:
                img_rects = page.get_image_rects()
            except Exception:
                img_rects = []
            if not img_rects:
                continue
            try:
                annots = list(page.annots()) if page.annots() else []
            except Exception:
                annots = []
            for annot in annots:
                info = annot.info
                desc = info.get("content", "") or info.get("alt", "") or ""
                if desc:
                    key = f"page{page_idx + 1}_{annot.rect}"
                    alt[key] = desc
        return alt

    @staticmethod
    def _text_in_rect(page: fitz.Page, rect: fitz.Rect) -> str:
        """Return text that overlaps *rect* on *page*."""
        try:
            words = page.get_text("words")
        except Exception:
            return ""
        parts: List[str] = []
        for w in words:
            wx0, wy0, wx1, wy1 = w[0], w[1], w[2], w[3]
            overlap = not (
                wx1 < rect.x0 or wx0 > rect.x1
                or wy1 < rect.y0 or wy0 > rect.y1
            )
            if overlap:
                parts.append(w[4])
        return " ".join(parts)


# ---------------------------------------------------------------------------
# Binding
# ---------------------------------------------------------------------------

class MetadataBinder:
    """Binds extracted PDF metadata onto a python-pptx Presentation."""

    def bind_metadata(
        self,
        prs: Presentation,
        pdf_metadata: PDFMetadata,
        page_mapping: Dict[int, int],
    ) -> None:
        if pdf_metadata.bookmarks:
            self._bind_bookmarks(prs, pdf_metadata.bookmarks)

        if pdf_metadata.hyperlinks:
            self._bind_hyperlinks(prs, pdf_metadata.hyperlinks, page_mapping)

        if pdf_metadata.alt_texts:
            self._bind_alt_texts(prs, pdf_metadata.alt_texts, page_mapping)

    # -- hyperlinks ---------------------------------------------------------

    def _bind_hyperlinks(
        self,
        prs: Presentation,
        hyperlinks: List[Hyperlink],
        page_mapping: Dict[int, int],
    ) -> None:
        for hl in hyperlinks:
            slide_idx = page_mapping.get(hl.page)
            if slide_idx is None or slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            self._apply_hyperlink_to_slide(slide, hl)

    def _apply_hyperlink_to_slide(self, slide, hl: Hyperlink) -> None:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_left = shape.left if shape.left is not None else 0
            shape_top = shape.top if shape.top is not None else 0
            shape_right = shape_left + (shape.width or 0)
            shape_bottom = shape_top + (shape.height or 0)
            if not self._bbox_overlap(
                (shape_left, shape_top, shape_right, shape_bottom),
                hl.bbox,
            ):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if self._bbox_overlap(
                        (shape_left, shape_top, shape_right, shape_bottom),
                        hl.bbox,
                    ):
                        run.hyperlink.address = hl.uri
                        return

    # -- bookmarks ----------------------------------------------------------

    def _bind_bookmarks(
        self,
        prs: Presentation,
        bookmarks: List[Bookmark],
    ) -> None:
        subject_parts = []
        for bm in bookmarks:
            subject_parts.append(
                f"{'  ' * (bm.level - 1)}{bm.title} (p{bm.page})"
            )
        prs.core_properties.subject = "\n".join(subject_parts)
        self._add_toc_slide(prs, bookmarks)

    def _add_toc_slide(
        self,
        prs: Presentation,
        bookmarks: List[Bookmark],
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.title.text = "Table of Contents"
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                tf = shape.text_frame
                tf.clear()
                for bm in bookmarks:
                    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
                    p.text = f"{'  ' * (bm.level - 1)}{bm.title}"
                    p.level = bm.level - 1
                break
        toc_idx = len(prs.slides) - 1
        for i in range(toc_idx, 0, -1):
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i - 1])
            prs.slides._sldIdLst.append(prs.slides._sldIdLst[-1] if len(prs.slides._sldIdLst) > 1 else prs.slides._sldIdLst[0])
        self._move_slide(prs, toc_idx, 0)

    @staticmethod
    def _move_slide(prs: Presentation, old_idx: int, new_idx: int) -> None:
        sldIdLst = prs.slides._sldIdLst
        sldId = sldIdLst[old_idx]
        sldIdLst.remove(sldId)
        if new_idx >= len(sldIdLst):
            sldIdLst.append(sldId)
        else:
            sldIdLst.insert(new_idx, sldId)

    # -- alt-text ------------------------------------------------------------

    def _bind_alt_texts(
        self,
        prs: Presentation,
        alt_texts: Dict[str, str],
        page_mapping: Dict[int, int],
    ) -> None:
        for key, desc in alt_texts.items():
            parts = key.split("_", 1)
            if len(parts) < 2:
                continue
            try:
                pdf_page = int(parts[0].replace("page", ""))
            except ValueError:
                continue
            slide_idx = page_mapping.get(pdf_page)
            if slide_idx is None or slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    shape.alt_text = desc
                    break

    # -- helpers ------------------------------------------------------------

    @staticmethod
    def _bbox_overlap(a, b) -> bool:
        ax0, ay0, ax1, ay1 = a
        bx0, by0, bx1, by1 = b
        ax0, ay0, ax1, ay1 = float(ax0), float(ay0), float(ax1), float(ay1)
        bx0, by0, bx1, by1 = float(bx0), float(by0), float(bx1), float(by1)
        return not (ax1 < bx0 or ax0 > bx1 or ay1 < by0 or ay0 > by1)


# ---------------------------------------------------------------------------
# Standalone CLI
# ---------------------------------------------------------------------------

def _cli_dump(pdf_path: str) -> None:
    with MetadataExtractor(pdf_path) as ext:
        meta = ext.extract_all()

    print(f"PDF: {pdf_path}")
    print(f"\n--- Hyperlinks ({len(meta.hyperlinks)}) ---")
    for hl in meta.hyperlinks:
        print(f"  p{hl.page}  {hl.uri}  bbox={[round(v,1) for v in hl.bbox]}  text={hl.text!r}")

    print(f"\n--- Bookmarks ({len(meta.bookmarks)}) ---")
    for bm in meta.bookmarks:
        print(f"  {'  ' * (bm.level - 1)}{bm.title}  (page {bm.page})")

    print(f"\n--- Page labels ({len(meta.page_labels)}) ---")
    for pg, lbl in meta.page_labels.items():
        print(f"  page {pg}: {lbl}")

    print(f"\n--- Alt-texts ({len(meta.alt_texts)}) ---")
    for key, desc in meta.alt_texts.items():
        print(f"  {key}: {desc}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python metadata_layer.py <pdf_path>")
        sys.exit(1)
    _cli_dump(sys.argv[1])
