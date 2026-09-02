from pptx import Presentation
prs = Presentation('output_test.pptx')
print(f'Slide count: {len(prs.slides)}')
print(f'Slide size: {prs.slide_width} x {prs.slide_height} EMUs')
for i, slide in enumerate(prs.slides):
    print(f'\nSlide {i+1}: {len(slide.shapes)} shapes')
    for shape in slide.shapes:
        text = ''
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text[:80]
        print(f'  Shape type={shape.shape_type} pos=({shape.left},{shape.top}) size=({shape.width}x{shape.height})')
        if text:
            print(f'    Text: "{text}"')
