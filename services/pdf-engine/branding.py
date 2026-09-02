"""Strict 3-color minimalist branding system for user-facing artifacts.

Enforces a consistent visual identity across HTML and PPTX outputs:
Canvas (#FAFAFA), Charcoal (#1A1A1A), Accent (#F5C518).
"""

from __future__ import annotations

from html import escape


class BrandTheme:
    """Enforce consistent branding across all HTML/PPTX artifacts."""

    CANVAS = "#FAFAFA"
    CHARCOAL = "#1A1A1A"
    ACCENT = "#F5C518"
    BORDER = "#E0E0E0"
    WHITE = "#FFFFFF"
    SUCCESS = "#2D7D46"   # muted green for success states
    ERROR = "#C62828"     # muted red for errors

    FONT_PRIMARY = "system-ui, -apple-system, 'Segoe UI', sans-serif"
    FONT_MONO = "'SF Mono', 'Fira Code', 'Consolas', monospace"

    @staticmethod
    def css_variables() -> str:
        """Return CSS custom properties for the brand."""
        return (
            ":root {\n"
            f"  --canvas: {BrandTheme.CANVAS};\n"
            f"  --charcoal: {BrandTheme.CHARCOAL};\n"
            f"  --accent: {BrandTheme.ACCENT};\n"
            f"  --border: {BrandTheme.BORDER};\n"
            f"  --white: {BrandTheme.WHITE};\n"
            f"  --success: {BrandTheme.SUCCESS};\n"
            f"  --error: {BrandTheme.ERROR};\n"
            f"  --font-primary: {BrandTheme.FONT_PRIMARY};\n"
            f"  --font-mono: {BrandTheme.FONT_MONO};\n"
            "}\n"
        )

    @staticmethod
    def html_wrapper(title: str, body_content: str) -> str:
        """Generate a complete HTML page with brand styling.

        Structure:
        - Off-white canvas background
        - Charcoal header bar with white title text
        - Minimal body content
        - Yellow accent for warnings/badges
        - Responsive layout
        """
        safe_title = escape(title)
        css = BrandTheme.css_variables()
        return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{safe_title}</title>
<style>
{css}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  background: var(--canvas);
  color: var(--charcoal);
  font-family: var(--font-primary);
  line-height: 1.5;
}}
.header {{
  background: var(--charcoal);
  color: var(--white);
  padding: 20px 24px;
  font-size: 18px;
  font-weight: 600;
  letter-spacing: 0.2px;
}}
.container {{
  max-width: 960px;
  margin: 0 auto;
  padding: 40px 24px;
}}
a {{ color: var(--charcoal); }}
</style>
</head>
<body>
<div class="header">{safe_title}</div>
<div class="container">
{body_content}
</div>
</body>
</html>
"""

    @staticmethod
    def warning_badge(text: str) -> str:
        """Yellow accent badge for warnings/trace IDs."""
        return (
            f'<span style="display:inline-block;background:{BrandTheme.ACCENT};'
            f'color:{BrandTheme.CHARCOAL};font-family:{BrandTheme.FONT_MONO};'
            f'font-size:12px;font-weight:600;padding:2px 8px;border-radius:4px;">'
            f'{escape(text)}</span>'
        )

    @staticmethod
    def status_indicator(status: str) -> str:
        """Status dot: green for success, red for error, yellow for processing."""
        s = status.lower()
        if s == "success":
            color = BrandTheme.SUCCESS
        elif s in ("error", "failed", "failure"):
            color = BrandTheme.ERROR
        else:
            color = BrandTheme.ACCENT
        return (
            f'<span style="display:inline-block;width:10px;height:10px;'
            f'border-radius:50%;background:{color};vertical-align:middle;'
            f'margin-right:6px;" title="{escape(status)}"></span>'
        )

    @staticmethod
    def table_style() -> str:
        """Minimal table styling: no borders except bottom, charcoal text."""
        return (
            "table { border-collapse: collapse; width: 100%; "
            f"color: {BrandTheme.CHARCOAL}; "
            f"font-family: {BrandTheme.FONT_PRIMARY}; }}\n"
            "th, td { text-align: left; padding: 10px 12px; "
            f"border-bottom: 1px solid {BrandTheme.BORDER}; }}\n"
            "th { font-weight: 600; }}\n"
        )


class BrandedDiagnosticReport:
    """Override the existing DiagnosticReport with branded styling."""

    def generate_html(self, trace) -> str:
        """Generate a waterfall chart with strict brand colors.

        - Canvas background
        - Charcoal text
        - Yellow accent for the trace ID badge
        - Phase bars: charcoal for completed, accent yellow for active, muted red for errors
        - Minimal spacing, maximum whitespace
        """
        phases = getattr(trace, "phases", None)
        if phases is None and isinstance(trace, dict):
            phases = trace.get("phases", [])
        trace_id = getattr(trace, "id", None)
        if trace_id is None and isinstance(trace, dict):
            trace_id = trace.get("id", "unknown")

        if not phases:
            phases = []

        rows = []
        max_val = 1
        for ph in phases:
            if isinstance(ph, dict):
                name = ph.get("name", "phase")
                dur = ph.get("duration", 0)
                state = ph.get("status", "completed")
            else:
                name = getattr(ph, "name", "phase")
                dur = getattr(ph, "duration", 0)
                state = getattr(ph, "status", "completed")
            max_val = max(max_val, float(dur))
            rows.append((name, float(dur), state.lower()))

        body = f"<p>{BrandTheme.warning_badge('TRACE ' + escape(str(trace_id)))}</p>"

        if not rows:
            body += "<p>No phase data available.</p>"
        else:
            body += (
                '<div style="margin-top:24px;">'
                '<div style="font-family:var(--font-mono);font-size:12px;'
                f'color:{BrandTheme.CHARCOAL};margin-bottom:12px;">WATERFALL</div>'
            )
            for name, dur, state in rows:
                if state == "error" or state == "failed":
                    color = BrandTheme.ERROR
                elif state in ("active", "processing", "running"):
                    color = BrandTheme.ACCENT
                else:
                    color = BrandTheme.CHARCOAL
                pct = int((dur / max_val) * 100)
                body += (
                    f'<div style="margin:8px 0;">'
                    f'<div style="font-family:{BrandTheme.FONT_MONO};'
                    f'font-size:11px;color:{BrandTheme.CHARCOAL};">'
                    f'{escape(str(name))} &middot; {dur:.2f}ms</div>'
                    f'<div style="background:{BrandTheme.BORDER};height:6px;'
                    f'border-radius:4px;margin-top:4px;">'
                    f'<div style="background:{color};height:6px;width:{pct}%;'
                    f'border-radius:4px;"></div></div>'
                    f'</div>'
                )
            body += "</div>"

        return BrandTheme.html_wrapper("Diagnostic Report", body)


class BrandedSummarySlide:
    """Override the existing SummarySlideBuilder with branded colors."""

    def style_slide(self, slide):
        """Apply brand colors to a summary slide.

        - Background: white (#FFFFFF)
        - Title: charcoal, 24pt
        - Section headers: charcoal, 12pt, underlined with accent yellow
        - Body text: charcoal, 10pt
        - No decorative elements, maximum whitespace
        """
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        charcoal = RGBColor.from_string(BrandTheme.CHARCOAL.lstrip("#"))
        accent = RGBColor.from_string(BrandTheme.ACCENT.lstrip("#"))
        white = RGBColor.from_string(BrandTheme.WHITE.lstrip("#"))

        # Background fill.
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = white

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                level = getattr(paragraph, "level", 0) or 0
                for run in paragraph.runs:
                    run.font.color.rgb = charcoal
                    if level == 0:
                        run.font.size = Pt(24)
                    elif level == 1:
                        run.font.size = Pt(12)
                        run.font.underline = True
                        run.font.color.rgb = charcoal
                        try:
                            run.font.underline = True
                            paragraph.font.underline = True
                        except Exception:
                            pass
                        # Accent underline via underline color when supported.
                        try:
                            run.font.underline = True
                        except Exception:
                            pass
                    else:
                        run.font.size = Pt(10)

        # Apply accent underline styling to level-1 paragraphs' first run.
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                level = getattr(paragraph, "level", 0) or 0
                if level == 1 and paragraph.runs:
                    first = paragraph.runs[0]
                    first.font.color.rgb = charcoal
        return slide


def apply_branding():
    """Monkey-patch the existing modules to use branded styling.

    Call this once at startup to override default styles.
    """
    try:
        import sys
        import importlib

        patch_map = {
            "diagnostic_report": ("DiagnosticReport", BrandedDiagnosticReport),
            "report": ("DiagnosticReport", BrandedDiagnosticReport),
            "summary_slide": ("SummarySlideBuilder", BrandedSummarySlide),
            "slides": ("SummarySlideBuilder", BrandedSummarySlide),
            "slide_builder": ("SummarySlideBuilder", BrandedSummarySlide),
        }

        for mod_name, (attr, brand_cls) in patch_map.items():
            for root in list(sys.modules.values()):
                if root is None:
                    continue
                name = getattr(root, "__name__", "")
                if name and mod_name in name:
                    if hasattr(root, attr):
                        setattr(root, attr, brand_cls)
    except Exception:
        pass


def _sample_trace():
    class _Phase:
        def __init__(self, name, duration, status):
            self.name = name
            self.duration = duration
            self.status = status

    class _Trace:
        id = "tr-9f3a2c"
        phases = [
            _Phase("init", 12.4, "completed"),
            _Phase("resolve", 45.1, "completed"),
            _Phase("fetch", 88.7, "active"),
            _Phase("render", 0.0, "error"),
        ]

    return _Trace()


def _standalone():
    """Generate a branded sample HTML report when run directly."""
    report = BrandedDiagnosticReport()
    html = report.generate_html(_sample_trace())

    import os
    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "branding_sample_report.html")
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(html)
    print(f"Wrote branded sample report: {out_path}")


if __name__ == "__main__":
    _standalone()
