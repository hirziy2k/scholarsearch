#!/usr/bin/env python3
"""Deterministic PPTX compiler — Brand Immutability Edict.

The python-pptx compiler is STRICTLY FORBIDDEN from applying programmatic styling.
It exclusively ingests the master .potx template and injects text into the
pre-existing, strictly branded bounding boxes defined by the master layouts.

Usage:
  python compile_pptx.py --doc <id> [--output output.pptx] [--template master_template.pptx]
"""

import json
import sys
import os
import sqlite3
import copy

from pptx import Presentation
from pptx.util import Inches, Pt

DB_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide_state.sqlite")
TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "master_template.potx")

# Layout mapping: layout_directive -> slide layout index in template
LAYOUT_MAP = {
    "title": 0,
    "content": 1,
    "comparison": 2,
    "process": 3,
    "definition": 4,
    "governance": 5,
}


def get_verified_slides(doc_id, db_path=DB_PATH):
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    rows = conn.execute(
        "SELECT slide_index, title, bullets, notes, source_page, status, "
        "layout_directive, visual_metadata "
        "FROM slides WHERE doc_id=? AND status IN ('done', 'override') "
        "ORDER BY slide_index",
        (doc_id,),
    ).fetchall()
    conn.close()

    slides = []
    for r in rows:
        try:
            bullets = json.loads(r["bullets"]) if isinstance(r["bullets"], str) else r["bullets"]
        except (json.JSONDecodeError, TypeError):
            bullets = []
        try:
            vm = json.loads(r["visual_metadata"]) if isinstance(r["visual_metadata"], str) else {}
        except (json.JSONDecodeError, TypeError):
            vm = {}
        slides.append({
            "slide_index": r["slide_index"],
            "title": r["title"],
            "bullets": bullets,
            "notes": r["notes"] or "",
            "source_page": r["source_page"],
            "status": r["status"],
            "layout_directive": r["layout_directive"] or "content",
            "visual_metadata": vm,
        })
    return slides


def fill_textbox(shape, text, append=False):
    """Inject text into an existing textbox shape. No styling applied."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if append:
        p = tf.add_paragraph()
        p.text = str(text)
    else:
        tf.paragraphs[0].text = str(text)


def inject_slide(prs, slide_idx, s, layout_idx):
    """Clone a template slide and inject content into bounding boxes."""
    template_slide = prs.slides[slide_idx]
    new_slide = prs.slides.add_slide(template_slide.slide_layout)

    # Copy background and shapes from template
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # Find placeholder shapes by their text content
    for shape in new_slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()

        if text == "[TITLE]":
            fill_textbox(shape, s["title"])
        elif text == "[BULLETS]":
            tf = shape.text_frame
            tf.paragraphs[0].text = ""
            for i, b in enumerate(s["bullets"]):
                if i == 0:
                    tf.paragraphs[0].text = str(b)
                else:
                    p = tf.add_paragraph()
                    p.text = str(b)
        elif text == "[NOTES]":
            fill_textbox(shape, s["notes"])
        elif text == "[SUBTITLE]":
            fill_textbox(shape, s["notes"] or f"Page {s['source_page']}")
        elif text == "[STEPS]":
            vm = s.get("visual_metadata", {})
            steps = vm.get("steps", s["bullets"])
            tf = shape.text_frame
            tf.paragraphs[0].text = ""
            for i, step in enumerate(steps):
                prefix = f"{i+1}. " if not str(step).startswith(tuple(str(j) for j in range(10))) else ""
                if i == 0:
                    tf.paragraphs[0].text = f"{prefix}{step}"
                else:
                    p = tf.add_paragraph()
                    p.text = f"{prefix}{step}"
        elif text == "[TERM]":
            vm = s.get("visual_metadata", {})
            fill_textbox(shape, vm.get("term", s["title"]))
        elif text == "[DEFINITION]":
            vm = s.get("visual_metadata", {})
            fill_textbox(shape, vm.get("definition", ". ".join(s["bullets"])))
        elif text == "[LEFT]":
            vm = s.get("visual_metadata", {})
            fill_textbox(shape, vm.get("left_label", ". ".join(s["bullets"][:2])))
        elif text == "[RIGHT]":
            vm = s.get("visual_metadata", {})
            fill_textbox(shape, vm.get("right_label", ". ".join(s["bullets"][2:])))
        elif text == "[STANDARD]":
            vm = s.get("visual_metadata", {})
            fill_textbox(shape, vm.get("standard", s["title"]))
        elif text == "[CRITERIA]":
            vm = s.get("visual_metadata", {})
            fill_textbox(shape, vm.get("criteria", ". ".join(s["bullets"][:2])))
        elif text == "[EVIDENCE]":
            vm = s.get("visual_metadata", {})
            fill_textbox(shape, vm.get("evidence", ". ".join(s["bullets"][2:])))

    # Set notes
    if s["notes"]:
        notes_slide = new_slide.notes_slide
        notes_slide.notes_text_frame.text = s["notes"]

    return new_slide


def compile_pptx(doc_id, output_path=None, template_path=TEMPLATE_PATH):
    """Compile slides to branded PPTX.

    Zero-Persistence Mandate: If output_path is None, returns a BytesIO
    buffer instead of writing to disk. The local directory remains sterile.
    """
    from io import BytesIO

    slides = get_verified_slides(doc_id)
    if not slides:
        return {"ok": False, "error": "No verified slides found"}

    if not os.path.exists(template_path):
        return {"ok": False, "error": f"Template not found: {template_path}"}

    # Load template as base
    prs = Presentation(template_path)
    template_slide_count = len(prs.slides)

    # Inject each verified slide
    for s in slides:
        layout_idx = LAYOUT_MAP.get(s["layout_directive"], 1)
        inject_slide(prs, layout_idx, s, layout_idx)

    if output_path is None:
        # Zero-Persistence: return buffer, no disk write
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return {
            "ok": True,
            "buffer": buf,
            "size": buf.getbuffer().nbytes,
            "slides_compiled": len(slides),
            "template_used": template_path,
            "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

    prs.save(output_path)
    return {
        "ok": True,
        "output": output_path,
        "slides_compiled": len(slides),
        "template_used": template_path,
    }


def main(argv):
    if "--doc" not in argv:
        print(__doc__)
        return 1
    doc_id = int(argv[argv.index("--doc") + 1])
    output = argv[argv.index("--output") + 1] if "--output" in argv else f"slides_doc{doc_id}.pptx"
    template = argv[argv.index("--template") + 1] if "--template" in argv else TEMPLATE_PATH

    result = compile_pptx(doc_id, output, template)
    print(json.dumps(result, indent=2))
    return 0 if result.get("ok") else 1


if __name__ == "__main__":
    sys.exit(main(sys.argv))
