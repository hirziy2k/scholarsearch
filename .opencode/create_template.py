#!/usr/bin/env python3
"""Create the master .potx template with branded visual hierarchy.

Brand Immutability Edict: This template hardcodes the exact visual hierarchy
at the file level. The compiler must never override these styles.

Colors:
  - Background: off-white #FAF9F6
  - Typography: charcoal #333
  - Accent: banana-yellow #FFD700
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors
BG_COLOR = RGBColor(0xFA, 0xF9, 0xF6)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0xFF, 0xD7, 0x00)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def create_template(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Layout 0: Title Slide
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, BG_COLOR)
    # Accent bar
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(2.8), Inches(2), Inches(0.08))  # thin rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    # Title placeholder
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.5),
                "[TITLE]", 36, CHARCOAL, bold=True)
    # Subtitle placeholder
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.0),
                "[SUBTITLE]", 20, GRAY)

    # Layout 1: Content (title + 4 bullets)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, BG_COLOR)
    add_textbox(slide2, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8),
                "[TITLE]", 28, CHARCOAL, bold=True)
    # Accent underline
    shape2 = slide2.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(3), Inches(0.06))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = ACCENT
    shape2.line.fill.background()
    # Bullets area
    add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0),
                "[BULLETS]", 18, GRAY)
    # Notes
    notes = slide2.notes_slide
    notes.notes_text_frame.text = "[NOTES]"

    # Layout 2: Comparison (left/right)
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3, BG_COLOR)
    add_textbox(slide3, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8),
                "[TITLE]", 28, CHARCOAL, bold=True)
    shape3 = slide3.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(3), Inches(0.06))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = ACCENT
    shape3.line.fill.background()
    # Left column
    add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
                "[LEFT]", 18, GRAY)
    # Vertical divider
    div = slide3.shapes.add_shape(1, Inches(6.5), Inches(1.8), Inches(0.06), Inches(4.5))
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT
    div.line.fill.background()
    # Right column
    add_textbox(slide3, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0),
                "[RIGHT]", 18, GRAY)

    # Layout 3: Process (numbered steps)
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide4, BG_COLOR)
    add_textbox(slide4, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8),
                "[TITLE]", 28, CHARCOAL, bold=True)
    shape4 = slide4.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(3), Inches(0.06))
    shape4.fill.solid()
    shape4.fill.fore_color.rgb = ACCENT
    shape4.line.fill.background()
    add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0),
                "[STEPS]", 18, GRAY)

    # Layout 4: Definition (term + definition)
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide5, BG_COLOR)
    # Term in accent color
    add_textbox(slide5, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2),
                "[TERM]", 32, ACCENT, bold=True)
    # Definition
    add_textbox(slide5, Inches(0.8), Inches(3.5), Inches(11.7), Inches(3.0),
                "[DEFINITION]", 20, CHARCOAL)

    # Layout 5: Governance (compliance: standard / criteria / evidence)
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide6, BG_COLOR)
    add_textbox(slide6, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8),
                "[TITLE]", 28, CHARCOAL, bold=True)
    shape6 = slide6.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(3), Inches(0.06))
    shape6.fill.solid()
    shape6.fill.fore_color.rgb = ACCENT
    shape6.line.fill.background()
    # Standard column
    add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(3.7), Inches(0.5),
                "Standard", 14, ACCENT, bold=True)
    add_textbox(slide6, Inches(0.8), Inches(2.3), Inches(3.7), Inches(4.0),
                "[STANDARD]", 16, CHARCOAL)
    # Criteria column
    add_textbox(slide6, Inches(4.8), Inches(1.8), Inches(3.7), Inches(0.5),
                "Criteria", 14, ACCENT, bold=True)
    add_textbox(slide6, Inches(4.8), Inches(2.3), Inches(3.7), Inches(4.0),
                "[CRITERIA]", 16, GRAY)
    # Evidence column
    add_textbox(slide6, Inches(8.8), Inches(1.8), Inches(3.7), Inches(0.5),
                "Evidence", 14, ACCENT, bold=True)
    add_textbox(slide6, Inches(8.8), Inches(2.3), Inches(3.7), Inches(4.0),
                "[EVIDENCE]", 16, GRAY)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    import os
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "master_template.potx")
    create_template(out)
    print(f"Template created: {out}")
