"""
Auto-Reflow Grouping for PDF-to-PPTX converter.

Groups text boxes that should vertically reflow together, so editing one
box pushes the others down instead of causing overlaps.
"""

from dataclasses import dataclass, field
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN


def _page_width_pt(prs: Presentation) -> float:
    return prs.slide_width / 12700.0


def _page_height_pt(prs: Presentation) -> float:
    return prs.slide_height / 12700.0


def _emu_to_pt(val: int) -> float:
    return val / 12700.0


def _pt_to_emu(val: float) -> int:
    return int(val * 12700.0)


def _left_edge(cluster: dict) -> float:
    """Return the left edge of a cluster in points."""
    return _emu_to_pt(cluster["left"])


def _top_edge(cluster: dict) -> float:
    return _emu_to_pt(cluster["top"])


def _bottom_edge(cluster: dict) -> float:
    return _emu_to_pt(cluster["top"] + cluster["height"])


def _height(cluster: dict) -> float:
    return _emu_to_pt(cluster["height"])


def _font_size(cluster: dict) -> float:
    return cluster.get("font_size", 14.0)


class VerticalStackDetector:
    """Detects vertical stacks of clusters that should reflow together."""

    def __init__(self, tolerance_pct: float = 0.10, gap_multiplier: float = 1.5):
        self.tolerance_pct = tolerance_pct
        self.gap_multiplier = gap_multiplier

    def detect_stacks(self, clusters: List[dict], page_height: float) -> List[List[int]]:
        if not clusters:
            return []

        n = len(clusters)
        # Build horizontal alignment groups first.
        left_edges = [_left_edge(c) for c in clusters]
        max_page_width = max(page_height, 1.0)  # use page_height as scale ref
        tolerance = max_page_width * self.tolerance_pct

        # Union-find for horizontal alignment
        parent = list(range(n))

        def find(x: int) -> int:
            while parent[x] != x:
                parent[x] = parent[parent[x]]
                x = parent[x]
            return x

        def union(a: int, b: int):
            ra, rb = find(a), find(b)
            if ra != rb:
                parent[ra] = rb

        for i in range(n):
            for j in range(i + 1, n):
                if abs(left_edges[i] - left_edges[j]) <= tolerance:
                    union(i, j)

        # Group by horizontal alignment
        h_groups: dict[int, List[int]] = {}
        for i in range(n):
            root = find(i)
            h_groups.setdefault(root, []).append(i)

        # Within each horizontal group, detect vertical stacks by gap consistency.
        stacks: List[List[int]] = []
        for members in h_groups.values():
            if len(members) == 1:
                continue
            # Sort top-to-bottom
            members_sorted = sorted(members, key=lambda i: _top_edge(clusters[i]))
            # Split into stacks based on gap consistency
            current_stack = [members_sorted[0]]
            for k in range(1, len(members_sorted)):
                prev_idx = members_sorted[k - 1]
                curr_idx = members_sorted[k]
                gap = _top_edge(clusters[curr_idx]) - _bottom_edge(clusters[prev_idx])
                avg_fs = (_font_size(clusters[prev_idx]) + _font_size(clusters[curr_idx])) / 2.0
                max_gap = avg_fs * self.gap_multiplier
                if 0 <= gap <= max_gap:
                    current_stack.append(curr_idx)
                else:
                    if len(current_stack) >= 2:
                        stacks.append(current_stack)
                    current_stack = [curr_idx]
            if len(current_stack) >= 2:
                stacks.append(current_stack)

        return stacks


@dataclass
class ReflowGroup:
    """A group of clusters that should vertically reflow together."""
    indices: List[int]
    anchor_index: int
    linked: bool = True

    def __repr__(self) -> str:
        return (
            f"ReflowGroup(indices={self.indices}, "
            f"anchor={self.anchor_index}, linked={self.linked})"
        )


class ReflowGroupBuilder:
    """Builds ReflowGroups from detected stacks; single clusters become standalone groups."""

    def build_groups(self, clusters: List[dict], stacks: List[List[int]]) -> List[ReflowGroup]:
        grouped_indices = set()
        groups: List[ReflowGroup] = []

        for stack in stacks:
            anchor = min(stack, key=lambda i: _top_edge(clusters[i]))
            groups.append(ReflowGroup(indices=list(stack), anchor_index=anchor, linked=True))
            grouped_indices.update(stack)

        for i in range(len(clusters)):
            if i not in grouped_indices:
                groups.append(ReflowGroup(indices=[i], anchor_index=i, linked=False))

        return groups


class PPTXReflowLinker:
    """Applies reflow properties and group shapes to a PPTX slide."""

    def apply_autofit(
        self,
        slide,
        groups: List[ReflowGroup],
        clusters: List[dict],
    ) -> None:
        """Set each text box to resize-to-fit with word wrap, and tag alt_text."""
        for group in groups:
            for idx in group.indices:
                shape = clusters[idx].get("shape")
                if shape is None or not hasattr(shape, "text_frame"):
                    continue
                tf = shape.text_frame
                tf.word_wrap = True
                # Auto-size to fit text vertically
                tf.auto_size = None  # will set via XML below

                # Use the XML API to set autofit to shrink-on-overflow
                from pptx.oxml.ns import qn
                bodyPr = tf._txBody.find(qn("a:bodyPr"))
                if bodyPr is not None:
                    bodyPr.set("wrap", "square")
                    # Remove any existing autofit elements
                    for child in list(bodyPr):
                        if child.tag.endswith("}spAutoFit") or child.tag.endswith("}normAutofit"):
                            bodyPr.remove(child)
                    # Add shrink-on-overflow autofit
                    auto_fit = bodyPr.makeelement(qn("a:spAutoFit"), {})
                    bodyPr.append(auto_fit)

                # Tag with group metadata in alt_text
                tag = (
                    f"reflow-group:{group.anchor_index}"
                    f"|linked:{group.linked}"
                    f"|members:{','.join(str(i) for i in group.indices)}"
                )
                existing_alt = ""
                if hasattr(shape, "alt_text") and shape.alt_text:
                    existing_alt = shape.alt_text
                if "reflow-group:" not in existing_alt:
                    shape.alt_text = f"{existing_alt} {tag}".strip() if existing_alt else tag

    def create_group_shape(
        self,
        slide,
        groups: List[ReflowGroup],
        clusters: List[dict],
        page_height: float,
        page_width: float,
    ) -> int:
        """Group stacked shapes into PPTX GroupShape for native reflow.
        Returns the number of groups created."""
        from pptx.oxml.ns import qn
        from lxml import etree

        grouped_count = 0

        for group in groups:
            if not group.linked or len(group.indices) < 2:
                continue

            shapes = []
            for idx in group.indices:
                shape = clusters[idx].get("shape")
                if shape is not None:
                    shapes.append(shape)

            if len(shapes) < 2:
                continue

            # Calculate bounding box
            left = min(_emu_to_pt(c["left"]) for c in [clusters[i] for i in group.indices])
            top = min(_emu_to_pt(c["top"]) for c in [clusters[i] for i in group.indices])
            right = max(
                _emu_to_pt(c["left"] + c["width"])
                for c in [clusters[i] for i in group.indices]
            )
            bottom = max(
                _emu_to_pt(c["top"] + c["height"])
                for c in [clusters[i] for i in group.indices]
            )

            # Create group shape via XML manipulation
            spTree = slide.shapes._spTree
            grpSp = spTree.makeelement(qn("p:grpSp"), {})
            grpSpPr = grpSp.makeelement(qn("p:grpSpPr"), {})

            # Set bounding box
            xfrm = grpSpPr.makeelement(qn("a:xfrm"), {})
            off = xfrm.makeelement(qn("a:off"), {"x": str(_pt_to_emu(left)), "y": str(_pt_to_emu(top))})
            ext = xfrm.makeelement(
                qn("a:ext"),
                {
                    "cx": str(_pt_to_emu(right - left)),
                    "cy": str(_pt_to_emu(bottom - top)),
                },
            )
            chOff = xfrm.makeelement(
                qn("a:chOff"), {"x": str(_pt_to_emu(left)), "y": str(_pt_to_emu(top))}
            )
            chExt = xfrm.makeelement(
                qn("a:chExt"),
                {
                    "cx": str(_pt_to_emu(right - left)),
                    "cy": str(_pt_to_emu(bottom - top)),
                },
            )
            xfrm.append(off)
            xfrm.append(ext)
            xfrm.append(chOff)
            xfrm.append(chExt)
            grpSpPr.append(xfrm)

            # NvGrpSpPr (required child)
            nvGrpSpPr = grpSp.makeelement(qn("p:nvGrpSpPr"), {})
            cNvPr = nvGrpSpPr.makeelement(qn("p:cNvPr"), {"id": "0", "name": f"ReflowGroup{grouped_count}"})
            cNvGrpSpPr = nvGrpSpPr.makeelement(qn("p:cNvGrpSpPr"), {})
            nvPr = nvGrpSpPr.makeelement(qn("p:nvPr"), {})
            nvGrpSpPr.append(cNvPr)
            nvGrpSpPr.append(cNvGrpSpPr)
            nvGrpSpPr.append(nvPr)

            grpSp.append(nvGrpSpPr)
            grpSp.append(grpSpPr)

            # Move shapes into the group element and remove from spTree
            for shape in shapes:
                sp_el = shape._element
                spTree.remove(sp_el)
                grpSp.append(sp_el)

            # Insert the group into the shape tree
            spTree.append(grpSp)
            grouped_count += 1

        return grouped_count


def apply_reflow_properties(
    slide,
    clusters: List[dict],
    page_height: float,
    page_width: float,
) -> int:
    """Detect stacks, build groups, apply autofit + grouping to the slide.

    Args:
        slide: A pptx slide object.
        clusters: List of dicts, each with keys:
            left, top, width, height (emu), font_size (pt, optional),
            shape (pptx shape, optional).
        page_height: Page height in points.
        page_width: Page width in points.

    Returns:
        Number of group shapes created.
    """
    detector = VerticalStackDetector()
    stacks = detector.detect_stacks(clusters, page_height)

    builder = ReflowGroupBuilder()
    groups = builder.build_groups(clusters, stacks)

    linker = PPTXReflowLinker()
    linker.apply_autofit(slide, groups, clusters)
    num_groups = linker.create_group_shape(slide, groups, clusters, page_height, page_width)

    return num_groups


if __name__ == "__main__":
    # Standalone demo / smoke test
    prs = Presentation()
    prs.slide_width = _pt_to_emu(960)
    prs.slide_height = _pt_to_emu(540)
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    pw = 960.0
    ph = 540.0

    # Create some sample shapes
    shapes_data = []
    titles = ["Title Box", "Subtitle Box", "Body Box", "Footer Box"]
    font_sizes = [28.0, 20.0, 14.0, 10.0]
    for i, (title, fs) in enumerate(zip(titles, font_sizes)):
        left = _pt_to_emu(100)
        top = _pt_to_emu(60 + i * 80)
        width = _pt_to_emu(760)
        height = _pt_to_emu(50)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(fs)
        shapes_data.append({
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "font_size": fs,
            "shape": txBox,
        })

    num = apply_reflow_properties(slide, shapes_data, ph, pw)
    print(f"Created {num} group shape(s) from {len(shapes_data)} clusters.")

    # Also print detected stacks
    detector = VerticalStackDetector()
    stacks = detector.detect_stacks(shapes_data, ph)
    print(f"Detected {len(stacks)} stack(s): {stacks}")

    builder = ReflowGroupBuilder()
    groups = builder.build_groups(shapes_data, stacks)
    for g in groups:
        print(f"  {g}")

    out = "reflow_demo.pptx"
    prs.save(out)
    print(f"Saved demo to {out}")
